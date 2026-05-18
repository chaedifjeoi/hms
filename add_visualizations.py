"""Insert 3 visualization slides into the 24/7 무상 기술지원 웹사이트 발표용 deck.

Inserts:
  - After original slide 9 (4-Tier 응답)        → "체감 응답 속도" 사다리 차트
  - After original slide 11 (선주 진입)          → "사용 빈도 점프" 비교 차트
  - After original slide 14 (Before/After 표)    → "정량 기대 효과" 비교 차트

Updates page-number labels (X / N) across the deck.
"""

from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- HD현대 brand palette (mined from the existing deck) ----------------
NAVY        = RGBColor(0x00, 0x2B, 0x5C)
NAVY_DARK   = RGBColor(0x00, 0x1A, 0x3D)
ORANGE      = RGBColor(0xFF, 0x7A, 0x00)
LIGHT_BLUE  = RGBColor(0xE8, 0xEE, 0xF7)
GRAY_TEXT   = RGBColor(0x33, 0x33, 0x33)
GRAY_SUB    = RGBColor(0x70, 0x80, 0x95)
LIGHT_GRAY  = RGBColor(0xCB, 0xD2, 0xD9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0x96, 0xA0)
RED         = RGBColor(0xC0, 0x39, 0x2B)

KFONT = "Apple SD Gothic Neo"

# ---- helpers ------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=GRAY_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=KFONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def add_header(slide, page_num, total_pages, kicker, title):
    """Replicate header used across the deck."""
    # top-left navy block
    add_rect(slide, Inches(0), Inches(0), Inches(0.3), Inches(1.5), NAVY)
    # kicker
    add_text(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.3),
             kicker, size=11, bold=True, color=ORANGE)
    # title
    add_text(slide, Inches(0.5), Inches(0.6), Inches(12), Inches(0.7),
             title, size=26, bold=True, color=NAVY)
    # footer orange bar
    add_rect(slide, Inches(0), Inches(7.4), Inches(13.33), Inches(0.1), ORANGE)
    # footer text
    add_text(slide, Inches(0.5), Inches(7.0), Inches(8), Inches(0.3),
             "HD현대마린솔루션테크 · 24/7 무상 기술지원 웹사이트 구축 제안",
             size=9, color=LIGHT_GRAY)
    add_text(slide, Inches(11.5), Inches(7.0), Inches(1.6), Inches(0.3),
             f"{page_num} / {total_pages}",
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ---- Visualization 1: 체감 응답 속도 사다리 (after slide 9) ------------
def build_response_ladder(slide):
    """Horizontal bar ladder visualizing 4-Tier SLA times."""
    add_text(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "Tier가 올라갈수록 응답 시간이 늘어나지만, 80% 이상의 문의는 Tier 1~2에서 즉시 처리됨",
             size=12, color=GRAY_SUB)

    # axis
    chart_left = Inches(2.6)
    chart_top = Inches(2.1)
    chart_w = Inches(8.6)
    chart_h = Inches(3.6)

    # bars: (tier_label, sub_label, time_label, ratio, color, share)
    tiers = [
        ("Tier 1", "자동 봇 + 지식베이스",     "0초",        0.04, ORANGE,    "≈ 50%"),
        ("Tier 2", "당번 엔지니어 (해외지사)",  "15분",       0.18, NAVY,      "≈ 35%"),
        ("Tier 3", "도메인 전문가 풀",         "2시간",      0.50, TEAL,      "≈ 12%"),
        ("Tier 4", "본사 R&D 엔지니어",        "24시간",     1.00, NAVY_DARK, "≈ 3%"),
    ]

    bar_h = Inches(0.6)
    gap = Inches(0.30)
    total_h = len(tiers) * bar_h + (len(tiers) - 1) * gap
    start_y = chart_top + (chart_h - total_h) // 2

    # background grid line at end
    add_rect(slide, chart_left, start_y - Inches(0.15),
             Emu(1), total_h + Inches(0.30), LIGHT_GRAY)

    for i, (tier, role, tlabel, ratio, color, share) in enumerate(tiers):
        y = start_y + i * (bar_h + gap)
        # tier label (left)
        add_text(slide, Inches(0.5), y + Inches(0.05),
                 Inches(2.0), Inches(0.3), tier,
                 size=14, bold=True, color=NAVY)
        add_text(slide, Inches(0.5), y + Inches(0.32),
                 Inches(2.0), Inches(0.3), role,
                 size=9, color=GRAY_SUB)

        # bar (background track)
        add_rect(slide, chart_left, y + Inches(0.18),
                 chart_w, Inches(0.25), LIGHT_BLUE)
        # filled portion
        bar_w = Emu(int(chart_w * ratio))
        add_rect(slide, chart_left, y + Inches(0.18),
                 bar_w, Inches(0.25), color)

        # time label at end of bar
        add_text(slide, chart_left + bar_w + Inches(0.1),
                 y + Inches(0.10), Inches(1.4), Inches(0.4),
                 tlabel, size=16, bold=True, color=color)

        # share at far right
        add_text(slide, Inches(12.4), y + Inches(0.18),
                 Inches(0.85), Inches(0.3), share,
                 size=11, color=GRAY_SUB, align=PP_ALIGN.RIGHT)

    # x-axis label
    add_text(slide, chart_left, chart_top + chart_h + Inches(0.0),
             chart_w, Inches(0.3), "  ←  빠름                                                  느림  →",
             size=9, color=GRAY_SUB, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(12.4), chart_top + chart_h + Inches(0.0),
             Inches(0.85), Inches(0.3), "예상 비중",
             size=9, color=GRAY_SUB, align=PP_ALIGN.RIGHT)

    # bottom KPI strip
    kpi_y = Inches(6.1)
    kpis = [
        ("0초",     "자동 응답 비중",         "Tier 1 자동화 적용"),
        ("15분",    "체감 SLA",              "Tier 2 당번 응답"),
        ("85%",     "Tier 1~2 처리율",       "사람 개입 최소화"),
        ("24/7",    "선주 체감 가용성",       "사람 24h 대기 불필요"),
    ]
    box_w = Inches(3.05)
    for i, (val, label, sub) in enumerate(kpis):
        x = Inches(0.5) + i * (box_w + Inches(0.05))
        add_rect(slide, x, kpi_y, box_w, Inches(0.7), LIGHT_BLUE)
        add_rect(slide, x, kpi_y, Inches(0.06), Inches(0.7), ORANGE)
        add_text(slide, x + Inches(0.18), kpi_y + Inches(0.05),
                 Inches(1.2), Inches(0.4), val,
                 size=20, bold=True, color=NAVY)
        add_text(slide, x + Inches(1.4), kpi_y + Inches(0.07),
                 box_w - Inches(1.5), Inches(0.3), label,
                 size=11, bold=True, color=NAVY)
        add_text(slide, x + Inches(1.4), kpi_y + Inches(0.36),
                 box_w - Inches(1.5), Inches(0.3), sub,
                 size=9, color=GRAY_SUB)


# ---- Visualization 2: 사용 빈도 점프 (after slide 11) ------------------
def build_frequency_chart(slide):
    """Compare contact frequency: OASIS (수리 견적) vs 본 제안 (일상 기술 문의)."""
    add_text(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "수리 견적은 연 1~2회뿐. 일상 기술 문의는 주/월 단위 — 빈도 자체가 20~50배 다름.",
             size=12, color=GRAY_SUB)

    # Two-panel comparison
    panel_w = Inches(5.9)
    panel_h = Inches(3.4)
    panel_y = Inches(2.0)

    # ---- LEFT: 오아시스 (low frequency)
    px = Inches(0.5)
    add_rect(slide, px, panel_y, panel_w, panel_h, LIGHT_BLUE)
    add_rect(slide, px, panel_y, panel_w, Inches(0.45), GRAY_SUB)
    add_text(slide, px + Inches(0.25), panel_y + Inches(0.07),
             panel_w, Inches(0.3),
             "오아시스 모델  ·  수리 견적 중개",
             size=13, bold=True, color=WHITE)

    add_text(slide, px + Inches(0.3), panel_y + Inches(0.6),
             panel_w - Inches(0.5), Inches(0.3),
             "트리거: 도크/수리 발생 시점만",
             size=10, color=GRAY_SUB)

    # giant number
    add_text(slide, px + Inches(0.3), panel_y + Inches(1.0),
             panel_w - Inches(0.6), Inches(1.3),
             "1~2",
             size=88, bold=True, color=GRAY_SUB)
    add_text(slide, px + Inches(0.3), panel_y + Inches(2.4),
             panel_w - Inches(0.6), Inches(0.4),
             "회 / 년",
             size=18, bold=True, color=GRAY_SUB)

    # year strip dots (12 months, 1-2 highlighted)
    strip_x = px + Inches(0.3)
    strip_y = panel_y + Inches(2.95)
    cell = Inches(0.42)
    for m in range(12):
        active = m in (4, 9)  # only 2 events
        c = ORANGE if active else WHITE
        add_rect(slide, strip_x + m * cell, strip_y, Inches(0.36), Inches(0.18),
                 c, line=GRAY_SUB if not active else None)

    # ---- RIGHT: 본 제안 (high frequency)
    px2 = Inches(6.9)
    add_rect(slide, px2, panel_y, panel_w, panel_h, NAVY)
    add_rect(slide, px2, panel_y, panel_w, Inches(0.45), ORANGE)
    add_text(slide, px2 + Inches(0.25), panel_y + Inches(0.07),
             panel_w, Inches(0.3),
             "본 제안  ·  일상 기술 문의 흐름",
             size=13, bold=True, color=WHITE)

    add_text(slide, px2 + Inches(0.3), panel_y + Inches(0.6),
             panel_w - Inches(0.5), Inches(0.3),
             "트리거: 매주~매월 발생하는 운영 이슈",
             size=10, color=LIGHT_GRAY)

    add_text(slide, px2 + Inches(0.3), panel_y + Inches(1.0),
             panel_w - Inches(0.6), Inches(1.3),
             "24~52",
             size=88, bold=True, color=ORANGE)
    add_text(slide, px2 + Inches(0.3), panel_y + Inches(2.4),
             panel_w - Inches(0.6), Inches(0.4),
             "회 / 년 (주 1회 ~ 월 2회 가정)",
             size=14, bold=True, color=WHITE)

    # year strip with majority dots highlighted
    strip2_x = px2 + Inches(0.3)
    strip2_y = panel_y + Inches(2.95)
    # roughly 24-30 events distributed in 12 months → fill many cells
    fill_pattern = [True, True, False, True, True, True, False, True, True, True, True, False]
    for m in range(12):
        c = ORANGE if fill_pattern[m] else NAVY_DARK
        add_rect(slide, strip2_x + m * cell, strip2_y, Inches(0.36), Inches(0.18), c,
                 line=LIGHT_GRAY if not fill_pattern[m] else None)

    # ---- Big arrow + multiplier
    arrow_x = Inches(6.4)
    arrow_y = Inches(5.7)
    arrow_w = Inches(0.5)
    arrow_h = Inches(0.5)
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.7))
    arr.fill.solid()
    arr.fill.fore_color.rgb = ORANGE
    arr.line.fill.background()
    tf = arr.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = ""
    r = p.add_run()
    r.text = "사용 빈도 ×20 ~ ×50  →  관계 자산화 가능 영역으로 이동"
    r.font.name = KFONT
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # legend
    add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "■ 접점 발생 (월별 시각화)   |   회색=미발생, 주황=발생   ·   년 단위 패턴 시뮬레이션",
             size=9, color=GRAY_SUB, align=PP_ALIGN.CENTER)


# ---- Visualization 3: Before / After 정량 효과 (after slide 14) ---------
def build_before_after_chart(slide):
    """Comparative bar chart for key Before/After metrics."""
    add_text(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "Before(메일·전화) → After(웹+24/7 응답) — 핵심 지표의 정량 점프",
             size=12, color=GRAY_SUB)

    metrics = [
        # (label, before_text, before_pct, after_text, after_pct, unit)
        ("기술 문의 평균 응답시간",   "12~24시간", 1.00, "15분",         0.02),
        ("케이스 디지털 자산화율",    "<5%",       0.05, "≥ 95%",        0.95),
        ("선주별 이력 일원화",        "메일함 산재", 0.10, "선박×년 단위", 1.00),
        ("야간/주말 1차 응답",        "X (다음 영업일)", 0.10, "○ 자동+당번", 1.00),
        ("FAQ 자동 해소 비중",        "0%",        0.00, "≈ 50%",        0.50),
        ("LCA 전환 가능 트리거 수",   "연 1~2회",  0.05, "주~월 단위",   0.85),
    ]

    chart_top = Inches(1.95)
    row_h = Inches(0.62)
    gap = Inches(0.10)
    label_w = Inches(3.6)
    bar_zone_x = Inches(4.3)
    bar_zone_w = Inches(6.3)
    text_zone_x = bar_zone_x + bar_zone_w + Inches(0.1)
    text_zone_w = Inches(2.4)

    # legend
    lg_y = Inches(1.45)
    add_rect(slide, Inches(9.9), lg_y, Inches(0.3), Inches(0.18), GRAY_SUB)
    add_text(slide, Inches(10.25), lg_y - Inches(0.02), Inches(1.0), Inches(0.25),
             "Before", size=10, color=GRAY_TEXT)
    add_rect(slide, Inches(11.3), lg_y, Inches(0.3), Inches(0.18), ORANGE)
    add_text(slide, Inches(11.65), lg_y - Inches(0.02), Inches(1.0), Inches(0.25),
             "After", size=10, color=GRAY_TEXT)

    for i, (label, b_text, b_pct, a_text, a_pct) in enumerate(metrics):
        y = chart_top + i * (row_h + gap)
        # row stripe (alternating bg)
        if i % 2 == 0:
            add_rect(slide, Inches(0.5), y - Inches(0.04),
                     Inches(12.3), row_h + Inches(0.08), LIGHT_BLUE)

        add_text(slide, Inches(0.6), y + Inches(0.16),
                 label_w, Inches(0.3), label,
                 size=11, bold=True, color=NAVY)

        # before bar (fills up to ~half of bar zone for the largest before value)
        before_w = Emu(int(bar_zone_w * 0.45 * b_pct))
        add_rect(slide, bar_zone_x, y + Inches(0.05),
                 max(before_w, Emu(1)), Inches(0.20),
                 GRAY_SUB)

        # after bar (fills full bar zone for max value)
        after_w = Emu(int(bar_zone_w * a_pct))
        add_rect(slide, bar_zone_x, y + Inches(0.30),
                 max(after_w, Emu(1)), Inches(0.22),
                 ORANGE)

        # both numeric labels go in the fixed text zone on the right
        add_text(slide, text_zone_x, y + Inches(0.0),
                 text_zone_w, Inches(0.25), b_text,
                 size=10, color=GRAY_SUB)
        add_text(slide, text_zone_x, y + Inches(0.28),
                 text_zone_w, Inches(0.25), a_text,
                 size=10, bold=True, color=ORANGE)

    # bottom callout
    cb_y = Inches(6.1)
    add_rect(slide, Inches(0.5), cb_y, Inches(12.3), Inches(0.7), NAVY)
    add_rect(slide, Inches(0.5), cb_y, Inches(0.12), Inches(0.7), ORANGE)
    add_text(slide, Inches(0.85), cb_y + Inches(0.07),
             Inches(11.5), Inches(0.3), "운영 효과 한 줄 요약",
             size=10, bold=True, color=ORANGE)
    add_text(slide, Inches(0.85), cb_y + Inches(0.30),
             Inches(11.5), Inches(0.4),
             "응답시간 ↓ 99%   ·   이력 자산화율 ↑ 19배   ·   야간/주말 공백 제거   ·   유상 전환 트리거 빈도 ↑ 20~50배",
             size=14, bold=True, color=WHITE)


# ---- slide insertion via XML (python-pptx lacks reorder API) ------------
def insert_blank_slide_at(prs, index):
    """Append a blank slide, then move it to `index` in the slide order."""
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    # find the truly blank layout if available
    for l in prs.slide_layouts:
        if l.name.lower() in ("blank",):
            layout = l
            break
    slide = prs.slides.add_slide(layout)
    # Move XML element to desired position
    sldIdLst = prs.slides._sldIdLst  # CT_SlideIdList
    children = list(sldIdLst)
    new_id = children[-1]
    sldIdLst.remove(new_id)
    sldIdLst.insert(index, new_id)
    return slide


def update_page_numbers(prs):
    """Rewrite the bottom-right 'X / N' label on every slide.

    Position-based to avoid false matches like a literal "24/7" KPI value.
    """
    total = len(prs.slides)
    import re
    pat = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")
    # Footer page-number box is anchored at top≈7.0", left≈11.5".
    MIN_TOP_EMU = Inches(6.8)
    MIN_LEFT_EMU = Inches(11.0)
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.top is None or shape.left is None:
                continue
            if shape.top < MIN_TOP_EMU or shape.left < MIN_LEFT_EMU:
                continue
            txt = shape.text_frame.text.strip()
            if pat.match(txt):
                # rewrite preserving formatting of the first run
                tf = shape.text_frame
                p = tf.paragraphs[0]
                # collect formatting from first run
                if p.runs:
                    first_run = p.runs[0]
                    fname = first_run.font.name
                    fsize = first_run.font.size
                    try:
                        fcolor = first_run.font.color.rgb
                    except Exception:
                        fcolor = LIGHT_GRAY
                else:
                    fname, fsize, fcolor = KFONT, Pt(10), LIGHT_GRAY
                # clear and rewrite
                for para in list(tf.paragraphs):
                    para_xml = para._p
                    para_xml.getparent().remove(para_xml)
                tf._txBody.append(etree.SubElement(tf._txBody, qn('a:p')))
                p2 = tf.paragraphs[0]
                p2.alignment = PP_ALIGN.RIGHT
                r = p2.add_run()
                r.text = f"{i} / {total}"
                if fname: r.font.name = fname
                if fsize: r.font.size = fsize
                if fcolor:
                    r.font.color.rgb = fcolor


def main():
    src = Path("/Users/sinchaeyeon/Desktop/채린/24_7_무상_기술지원_웹사이트_구축_제안_발표용.pptx")
    out = src  # overwrite
    prs = Presentation(str(src))

    # Insert positions are 0-indexed. We insert AFTER slide 9 (idx 9 = pos 9 after orig idx 8),
    # AFTER slide 11 (orig idx 10), AFTER slide 14 (orig idx 13).
    # Insert in reverse order so earlier indices stay stable.

    # 1) After original slide 14 (index 13 → insert at pos 14)
    s_after14 = insert_blank_slide_at(prs, 14)
    # 2) After original slide 11 (index 10 → insert at pos 11)
    s_after11 = insert_blank_slide_at(prs, 11)
    # 3) After original slide 9 (index 8 → insert at pos 9)
    s_after9 = insert_blank_slide_at(prs, 9)

    # Now total slides = 23. New slide ordering:
    #   1..9 unchanged
    #   10 = response ladder (new)
    #   11 = orig 10
    #   12 = orig 11
    #   13 = freq jump (new)
    #   14 = orig 12
    #   15 = orig 13
    #   16 = orig 14
    #   17 = before/after chart (new)
    #   18.. = orig 15..20

    total = len(prs.slides)
    add_header(s_after9,  10, total, "본론 3 · Q3 — 체감 응답 속도", "Tier별 SLA 사다리 — 즉시→24시간")
    build_response_ladder(s_after9)

    add_header(s_after11, 13, total, "본론 4 · Q4 — 사용 빈도 점프", "오아시스 대비 ×20~×50 — 관계 자산화의 동력")
    build_frequency_chart(s_after11)

    add_header(s_after14, 17, total, "결론 — 정량 기대 효과", "Before vs After 핵심 지표 비교")
    build_before_after_chart(s_after14)

    # Update all page-number labels
    update_page_numbers(prs)

    prs.save(str(out))
    print(f"✓ saved {out} ({total} slides)")


if __name__ == "__main__":
    main()
