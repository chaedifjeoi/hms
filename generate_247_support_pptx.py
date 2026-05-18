# -*- coding: utf-8 -*-
"""
24/7 무상 기술 지원 웹사이트 구축 제안 — PPT 생성
디자인 시스템 v2 준수: 맑은 고딕 단일 / 16:9 고정 / 헤더 좌표 고정 / 본문 고밀도
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ========== 디자인 시스템 v2 상수 ==========

FONT_NAME = '맑은 고딕'

# 컬러
INK = '111827'
TEXT = '374151'
MUTED = '6B7280'
LIGHT_MUTED = '9CA3AF'
BG_WHITE = 'FFFFFF'
BG_SOFT = 'F9FAFB'
BG_BOX = 'F3F4F6'
ACCENT = '1F4E79'
ACCENT_2 = '3A6FA0'
ACCENT_3 = '5C8FBE'
ACCENT_4 = '9DBBD5'
WARN = 'C0392B'
LINE = 'D9D9D9'
LINE_SOFT = 'E5E7EB'
DARK_BG = '1F2937'

# 슬라이드 좌표 (Inches)
SLIDE_W = 13.333
SLIDE_H = 7.5
M_L = 0.6      # 좌측 여백
M_R = 0.6      # 우측 여백
CONTENT_W = SLIDE_W - M_L - M_R   # 12.133
CHAPTER_Y = 0.5
CHAPTER_H = 0.3
TITLE_Y = 0.85
TITLE_H = 0.7
SUBTITLE_Y = 1.6
SUBTITLE_H = 0.45
DIVIDER_Y = 2.1
BODY_Y = 2.25
BODY_H = 4.8
FOOTER_Y = 7.05
FOOTER_H = 0.4


# ========== 폰트 헬퍼 (rFonts 강제 설정) ==========

def apply_font(run, size=11, bold=False, color=INK, italic=False):
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor.from_string(color)
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn('a:rFonts'))
    if rFonts is None:
        rFonts = etree.SubElement(rPr, qn('a:rFonts'))
    rFonts.set('eastAsia', FONT_NAME)
    rFonts.set('ascii', FONT_NAME)
    rFonts.set('hAnsi', FONT_NAME)


def set_para(para, align=PP_ALIGN.LEFT, space_before=0, space_after=2):
    para.alignment = align
    para.space_before = Pt(space_before)
    para.space_after = Pt(space_after)


# ========== 도형 헬퍼 ==========

def add_rect(slide, x, y, w, h, fill=BG_WHITE, line=None, line_w=0.5, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = RGBColor.from_string(line)
        shape.line.width = Pt(line_w)
    if not shadow:
        # 그림자 제거
        sppr = shape.shadow._element
        for sp in shape._element.findall('.//' + qn('a:effectLst')):
            sp.getparent().remove(sp)
        # 빈 effectLst 추가해서 상속 방지
        spPr = shape._element.spPr
        existing = spPr.find(qn('a:effectLst'))
        if existing is not None:
            spPr.remove(existing)
        etree.SubElement(spPr, qn('a:effectLst'))
    shape.text_frame.margin_left = Inches(0.08)
    shape.text_frame.margin_right = Inches(0.08)
    shape.text_frame.margin_top = Inches(0.05)
    shape.text_frame.margin_bottom = Inches(0.05)
    return shape


def add_line(slide, x, y, w, h, color=LINE, weight=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    spPr = shape._element.spPr
    existing = spPr.find(qn('a:effectLst'))
    if existing is not None:
        spPr.remove(existing)
    etree.SubElement(spPr, qn('a:effectLst'))
    return shape


def add_text(slide, x, y, w, h, text, size=11, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25,
             italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor

    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        apply_font(run, size=size, bold=bold, color=color, italic=italic)
    return tb


def add_text_in_shape(shape, text, size=11, bold=False, color=INK,
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        apply_font(run, size=size, bold=bold, color=color)


# ========== 헤더/푸터 (모든 슬라이드 동일 좌표) ==========

def add_header(slide, chapter_label, title, subtitle, page_num, total_pages,
               doc_title='24/7 무상 기술 지원 웹사이트'):
    # 챕터 라벨
    add_text(slide, M_L, CHAPTER_Y, CONTENT_W, CHAPTER_H,
             chapter_label, size=9, bold=True, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    # 제목
    add_text(slide, M_L, TITLE_Y, CONTENT_W, TITLE_H,
             title, size=24, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    # 부제목
    add_text(slide, M_L, SUBTITLE_Y, CONTENT_W, SUBTITLE_H,
             subtitle, size=14, bold=False, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    # 구분선
    add_line(slide, M_L, DIVIDER_Y, CONTENT_W, 0.015, color=LINE)
    # 푸터 좌
    add_text(slide, M_L, FOOTER_Y, CONTENT_W * 0.6, FOOTER_H,
             doc_title, size=9, color=LIGHT_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    # 푸터 우 (페이지)
    add_text(slide, M_L + CONTENT_W * 0.6, FOOTER_Y, CONTENT_W * 0.4, FOOTER_H,
             f'{page_num} / {total_pages}', size=9, color=LIGHT_MUTED,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


# ========== 카드 컴포넌트 ==========

def add_card(slide, x, y, w, h, title=None, body=None, kicker=None,
             accent_bar=False, fill=BG_WHITE, line=LINE_SOFT,
             title_size=12, body_size=10):
    """기본 카드: 제목 + 본문"""
    shape = add_rect(slide, x, y, w, h, fill=fill, line=line)
    if accent_bar:
        add_rect(slide, x, y, 0.06, h, fill=ACCENT, line=None)
    pad_x = 0.12 if not accent_bar else 0.18
    pad_y = 0.12
    cur_y = y + pad_y
    if kicker:
        add_text(slide, x + pad_x, cur_y, w - pad_x * 2, 0.22,
                 kicker, size=8, bold=True, color=ACCENT)
        cur_y += 0.22
    if title:
        add_text(slide, x + pad_x, cur_y, w - pad_x * 2, 0.32,
                 title, size=title_size, bold=True, color=INK)
        cur_y += 0.36
    if body:
        add_text(slide, x + pad_x, cur_y, w - pad_x * 2, h - (cur_y - y) - pad_y,
                 body, size=body_size, color=TEXT, line_spacing=1.3)


def add_kpi_card(slide, x, y, w, h, big, label, big_color=ACCENT, fill=BG_WHITE):
    add_rect(slide, x, y, w, h, fill=fill, line=LINE_SOFT)
    add_text(slide, x + 0.12, y + 0.08, w - 0.24, h * 0.55,
             big, size=28, bold=True, color=big_color, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + 0.12, y + h * 0.62, w - 0.24, h * 0.32,
             label, size=10, color=TEXT, line_spacing=1.25)


def add_takeaway(slide, text, y=None):
    """본문 하단 핵심 한 줄 박스"""
    if y is None:
        y = BODY_Y + BODY_H - 0.55
    add_rect(slide, M_L, y, CONTENT_W, 0.5, fill=BG_BOX, line=None)
    add_rect(slide, M_L, y, 0.05, 0.5, fill=ACCENT, line=None)
    add_text(slide, M_L + 0.18, y, CONTENT_W - 0.3, 0.5,
             text, size=12, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT, weight=2):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = RGBColor.from_string(color)
    shape.line.width = Pt(weight)
    # 화살표 머리
    ln = shape.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return shape


def add_section_bg(slide, color=BG_WHITE):
    # 슬라이드 전체 배경 사각형
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(color)
    bg.line.fill.background()
    # 가장 뒤로
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


# ========== 슬라이드 빌더들 ==========

TOTAL = 27


def slide_01_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)

    # 헤더 좌표 동일 적용
    add_text(s, M_L, CHAPTER_Y, CONTENT_W, CHAPTER_H,
             'TITLE', size=9, bold=True, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, M_L, TITLE_Y, CONTENT_W, TITLE_H,
             '24/7 무상 기술 지원 웹사이트 구축',
             size=24, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, M_L, SUBTITLE_Y, CONTENT_W, SUBTITLE_H,
             '선도기업처럼, 그러나 오아시스와 다르게',
             size=14, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    add_line(s, M_L, DIVIDER_Y, CONTENT_W, 0.015, color=LINE)

    # 본문: 좌측 발표자 카드 / 우측 KPI
    cw = CONTENT_W * 0.45
    add_card(s, M_L, BODY_Y + 0.3, cw, 2.0,
             kicker='발표 정보',
             title='HD현대마린솔루션테크',
             body='발표자  ·  AS 사업본부\n일자  ·  2026년 5월\n대상  ·  내부 전략 보고',
             accent_bar=True, body_size=11, title_size=14)

    # 우측 KPI 3종
    rx = M_L + cw + 0.3
    rw = CONTENT_W - cw - 0.3
    add_text(s, rx, BODY_Y + 0.3, rw, 0.3,
             '우리가 가진 것', size=10, bold=True, color=ACCENT)
    add_kpi_card(s, rx, BODY_Y + 0.7, rw / 3 - 0.1, 1.6,
                 '9,890', '척\nAM 솔루션 제공 선박')
    add_kpi_card(s, rx + rw / 3 + 0.05, BODY_Y + 0.7, rw / 3 - 0.1, 1.6,
                 '24/7', '시간\n스마트케어 센터')
    add_kpi_card(s, rx + 2 * (rw / 3) + 0.1, BODY_Y + 0.7, rw / 3 - 0.2, 1.6,
                 '350+', '척\n실시간 모니터링')

    # 하단 인용
    add_takeaway(s, '기술력은 충분하다 — 부족한 것은 고객이 닿는 디지털 창구 하나다',
                 y=BODY_Y + 3.2)

    add_text(s, M_L, FOOTER_Y, CONTENT_W * 0.6, FOOTER_H,
             '24/7 무상 기술 지원 웹사이트', size=9, color=LIGHT_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, M_L + CONTENT_W * 0.6, FOOTER_Y, CONTENT_W * 0.4, FOOTER_H,
             f'1 / {TOTAL}', size=9, color=LIGHT_MUTED,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


def slide_02_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'AGENDA', '오늘 다룰 7가지',
               '망한 사례에서 시작해 유상 전환까지', 2, TOTAL)

    items = [
        ('01', '오아시스가 망한 이유', '강제성 없는 중개의 한계'),
        ('02', '그럼에도 웹사이트', '선도기업은 이미 웹으로 갔다'),
        ('03', 'ISS 2.0 위에 탑재', '별도 사이트가 아니다'),
        ('04', '24/7 무상 기술 지원', '해외지사 시간대 릴레이'),
        ('05', 'PM 인수인계 자동화', 'AI가 정리하고 PM은 클로징한다'),
        ('06', '유상 전환 — 구독형', '라이프사이클 계약 3티어'),
        ('07', '기대효과', '정량 8가지 · 정성 3중 효과'),
        ('Q&A', '질의응답', ''),
    ]
    cols = 4
    rows = 2
    gap_x = 0.2
    gap_y = 0.2
    cw = (CONTENT_W - gap_x * (cols - 1)) / cols
    ch = (BODY_H - gap_y * (rows - 1) - 0.3) / rows

    for i, (num, title, desc) in enumerate(items):
        r = i // cols
        c = i % cols
        x = M_L + c * (cw + gap_x)
        y = BODY_Y + 0.15 + r * (ch + gap_y)
        is_qa = num == 'Q&A'
        add_card(s, x, y, cw, ch,
                 kicker=f'CHAPTER {num}' if not is_qa else 'Q&A',
                 title=title, body=desc,
                 accent_bar=True,
                 fill=BG_SOFT if is_qa else BG_WHITE,
                 title_size=13, body_size=10)


# ========== CH1. 오아시스가 망한 이유 ==========

def slide_03_oasis_what(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 01 · OASIS POSTMORTEM',
               '우리는 이미 한 번 실패했다',
               'HD현대마린솔루션 오아시스, 무엇을 했고 왜 멈췄나', 3, TOTAL)

    # 좌측 50% — 오아시스 6대 기능
    lw = CONTENT_W * 0.55
    add_text(s, M_L, BODY_Y + 0.1, lw, 0.3,
             '오아시스가 가졌던 6대 기능', size=11, bold=True, color=ACCENT)
    feats = [
        ('견적 비교', '수리업체 평점·견적 한 번에'),
        ('오아시스 보증', 'HGS 50년 무상보증 경험 기반'),
        ('실시간 대응', '커뮤니케이션을 웹으로'),
        ('중개 플랫폼', '고객사 ↔ 수리업체 매칭'),
        ('바우처 3~5%', '건당 할인 인센티브'),
        ('글로벌 포트', '해외 항구 수리 지원'),
    ]
    cols = 2
    cw = (lw - 0.15) / cols
    ch = 0.85
    for i, (t, d) in enumerate(feats):
        r = i // cols
        c = i % cols
        x = M_L + c * (cw + 0.15)
        y = BODY_Y + 0.5 + r * (ch + 0.12)
        add_card(s, x, y, cw, ch, title=t, body=d, title_size=11, body_size=9)

    # 우측 인용 박스
    rx = M_L + lw + 0.25
    rw = CONTENT_W - lw - 0.25
    add_rect(s, rx, BODY_Y + 0.1, rw, 3.5, fill=BG_BOX, line=None)
    add_rect(s, rx, BODY_Y + 0.1, 0.05, 3.5, fill=ACCENT, line=None)
    add_text(s, rx + 0.2, BODY_Y + 0.3, rw - 0.3, 0.3,
             '선주의 진짜 목소리', size=9, bold=True, color=ACCENT)
    add_text(s, rx + 0.2, BODY_Y + 0.65, rw - 0.3, 2.8,
             '"수리업체 견적 비교는 좋다.\n그런데 매번 로그인해서 들어가는 게 번거롭다."\n\n"메일 한 통이면 PM이 다 처리해주는데,\n굳이 웹사이트를 거쳐야 할 이유가 없다."\n\n"우리 배 정보는 우리 회사 시스템에 있지,\n중개 플랫폼에는 없다."',
             size=11, color=TEXT, line_spacing=1.45)

    add_takeaway(s, '기능은 풍부했지만, 선주에게 "쓰지 않으면 손해인 이유"가 없었다')


def slide_04_oasis_why_failed(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 01 · OASIS POSTMORTEM',
               '망한 진짜 이유 4가지',
               '기능이 아니라 진입 동기의 문제였다', 4, TOTAL)

    reasons = [
        ('01', '강제성 부재',
         '선주가 굳이 들어와야 할 이유가 없었다.\n메일 한 통이면 PM이 다 처리하는데,\n웹을 거쳐야 할 동기가 없다.'),
        ('02', '중개의 한계',
         '플랫폼이 데이터를 직접 소유하지 못했다.\n선주의 호선 이력은 OEM 시스템에 있고,\n중개 사이트에는 거래 기록만 남았다.'),
        ('03', '가격 경쟁 매몰',
         '핵심 메시지가 "비교 견적·할인"으로 굳었다.\nESG·품질·신뢰의 차별 메시지가 희석되어\n결국 단가 경쟁의 무대가 되었다.'),
        ('04', '운영 관성',
         'PM과 선주 모두 메일이 빠르고 익숙했다.\n새 도구가 기존 워크플로우보다\n명확히 빠르지 않으면 살아남지 못한다.'),
    ]
    cols = 2
    cw = (CONTENT_W - 0.25) / cols
    ch = 1.85
    for i, (num, t, d) in enumerate(reasons):
        r = i // cols
        c = i % cols
        x = M_L + c * (cw + 0.25)
        y = BODY_Y + 0.15 + r * (ch + 0.2)
        add_card(s, x, y, cw, ch,
                 kicker=f'REASON {num}', title=t, body=d,
                 accent_bar=True, title_size=14, body_size=10)

    add_takeaway(s, '4가지 모두 "기능 부족"이 아니라 "동기 부족" — 같은 실수를 반복하지 않는다')


def slide_05_oasis_lesson(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 01 · OASIS POSTMORTEM',
               '그래서 우리가 배운 것',
               '"쓸 수 있는" 사이트가 아니라 "쓰지 않으면 손해인" 사이트로', 5, TOTAL)

    # 가로 비교 다이어그램
    box_w = CONTENT_W * 0.42
    box_h = 1.8
    arrow_w = CONTENT_W - box_w * 2

    # 좌측 (오아시스)
    lx = M_L
    ly = BODY_Y + 0.4
    add_rect(s, lx, ly, box_w, box_h, fill=BG_BOX, line=LINE_SOFT)
    add_text(s, lx + 0.2, ly + 0.15, box_w - 0.4, 0.3,
             '과거 — 오아시스', size=10, bold=True, color=MUTED)
    add_text(s, lx + 0.2, ly + 0.5, box_w - 0.4, box_h - 0.6,
             '· 중개 비즈니스 (외부 플랫폼)\n· 가격 비교가 핵심 가치\n· 데이터는 거래 기록만\n· 가입·로그인 별도 필요\n· 강제 진입 동기 없음',
             size=11, color=TEXT, line_spacing=1.5)

    # 화살표
    add_arrow(s, lx + box_w + 0.15, ly + box_h / 2,
              lx + box_w + arrow_w - 0.15, ly + box_h / 2,
              color=ACCENT, weight=3)

    # 우측 (본 제안)
    rx = lx + box_w + arrow_w
    add_rect(s, rx, ly, box_w, box_h, fill=ACCENT, line=None)
    add_text(s, rx + 0.2, ly + 0.15, box_w - 0.4, 0.3,
             '본 제안 — 24/7 무상 기술 지원',
             size=10, bold=True, color=BG_WHITE)
    add_text(s, rx + 0.2, ly + 0.5, box_w - 0.4, box_h - 0.6,
             '· OEM AS 비즈니스 (자사 플랫폼)\n· 기술·이력·신뢰가 핵심 가치\n· 호선 데이터를 직접 소유\n· ISS 2.0 사용자 즉시 흡수\n· 24/7 무상 지원이 진입 미끼',
             size=11, color=BG_WHITE, line_spacing=1.5)

    add_takeaway(s, '기능을 더 넣는 게 아니라, 진입 동기를 다시 설계한다')


# ========== CH2. 그럼에도 웹사이트 ==========

def slide_06_benchmark(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 02 · WHY STILL WEB',
               '선도기업은 이미 웹으로 이동했다',
               'Wärtsilä Online · Kongsberg K-IMS', 6, TOTAL)

    cw = (CONTENT_W - 0.3) / 2
    ch = 3.4

    # Wärtsilä
    lx = M_L
    ly = BODY_Y + 0.15
    add_rect(s, lx, ly, cw, ch, fill=BG_WHITE, line=LINE_SOFT)
    add_rect(s, lx, ly, 0.06, ch, fill=ACCENT, line=None)
    add_text(s, lx + 0.2, ly + 0.15, cw - 0.3, 0.3,
             'WÄRTSILÄ ONLINE', size=9, bold=True, color=ACCENT)
    add_text(s, lx + 0.2, ly + 0.45, cw - 0.3, 0.4,
             '핀란드 · 글로벌 OEM', size=14, bold=True, color=INK)
    add_text(s, lx + 0.2, ly + 0.95, cw - 0.3, 2.3,
             '· TechRequest — 웹에서 기술 지원 요청\n· 클레임 / 워런티 등록 → 진행 추적\n· 견적 요청 → 상태 + 문서 조회\n· 매뉴얼 · 기술 회보 24/7 셀프서비스\n· 호선 정보를 고객이 직접 관리\n· 포털 내에서 엔지니어 직접 문의',
             size=11, color=TEXT, line_spacing=1.55)

    # Kongsberg
    rx = M_L + cw + 0.3
    add_rect(s, rx, ly, cw, ch, fill=BG_WHITE, line=LINE_SOFT)
    add_rect(s, rx, ly, 0.06, ch, fill=ACCENT, line=None)
    add_text(s, rx + 0.2, ly + 0.15, cw - 0.3, 0.3,
             'KONGSBERG K-IMS', size=9, bold=True, color=ACCENT)
    add_text(s, rx + 0.2, ly + 0.45, cw - 0.3, 0.4,
             '노르웨이 · 글로벌 OEM', size=14, bold=True, color=INK)
    add_text(s, rx + 0.2, ly + 0.95, cw - 0.3, 2.3,
             '· 선박 ↔ 육상 데이터 실시간 동기화\n· 글로벌 34개국 24/7 지원 네트워크\n· 지원·부품 요청 단일 폼 통합\n· 항해·운항·트림·배출 대시보드\n· 2026 KM Performance — 레거시 통합\n· 별도 앱 없이 브라우저로 접근',
             size=11, color=TEXT, line_spacing=1.55)

    # 하단 띠
    add_takeaway(s, 'AM 시장 2024 371억 USD → 2032 532억 USD · CAGR 4.6% — 디지털 창구는 이미 표준이다')


def slide_07_we_are_different(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 02 · WHY STILL WEB',
               '우리는 오아시스와 다르다',
               '중개사가 아니라 OEM AS사이고, 데이터가 우리에게 있다', 7, TOTAL)

    # 비교표 헤더
    tx = M_L
    ty = BODY_Y + 0.2
    th = 0.42
    col_label = CONTENT_W * 0.18
    col_oasis = CONTENT_W * 0.41
    col_us = CONTENT_W * 0.41

    # 헤더 행
    add_rect(s, tx, ty, col_label, th, fill=DARK_BG, line=None)
    add_rect(s, tx + col_label, ty, col_oasis, th, fill=DARK_BG, line=None)
    add_rect(s, tx + col_label + col_oasis, ty, col_us, th, fill=ACCENT, line=None)
    add_text(s, tx + 0.15, ty, col_label - 0.3, th,
             '항목', size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + col_label + 0.15, ty, col_oasis - 0.3, th,
             '오아시스 (과거)', size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + col_label + col_oasis + 0.15, ty, col_us - 0.3, th,
             '본 제안 (현재)', size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ('데이터', '외부 플랫폼 (거래 기록만)', '자사 시스템 (호선·이력 전체 보유)'),
        ('진입 경로', '별도 가입·로그인 필요', 'ISS 2.0 사용자 그대로 흡수'),
        ('핵심 가치', '가격 비교 · 견적 할인', '기술·이력·신뢰 · OEM 정통성'),
        ('강제성', '없음 — 안 써도 그만', '24/7 무상 기술 지원이 미끼'),
        ('수익 모델', '중개 수수료 (단발)', '구독형 LTSA (반복 매출)'),
    ]
    cur_y = ty + th
    rh = 0.55
    for i, (label, oasis, us) in enumerate(rows):
        bg = BG_SOFT if i % 2 == 0 else BG_WHITE
        add_rect(s, tx, cur_y, CONTENT_W, rh, fill=bg, line=None)
        add_line(s, tx, cur_y, CONTENT_W, 0.005, color=LINE_SOFT)
        add_text(s, tx + 0.15, cur_y, col_label - 0.3, rh,
                 label, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + col_label + 0.15, cur_y, col_oasis - 0.3, rh,
                 oasis, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + col_label + col_oasis + 0.15, cur_y, col_us - 0.3, rh,
                 us, size=10, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        cur_y += rh

    add_takeaway(s, '우리가 가진 것은 9,890척의 데이터다 — 오아시스는 가질 수 없었던 자산이다')


def slide_08_gap_essence(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 02 · WHY STILL WEB',
               'Gap의 본질',
               '기술력이 아니라 디지털 창구다', 8, TOTAL)

    # KPI 4개
    kpis = [
        ('9,890', '척\nAM 솔루션 제공 (2024 3Q)'),
        ('24/7', '시간\n부산 스마트케어 센터'),
        ('350+', '척\n실시간 모니터링'),
        ('5.3%', '연료 절감\nOceanWise AI 항로 최적화'),
    ]
    cw = (CONTENT_W - 0.3) / 4
    ch = 1.8
    for i, (big, label) in enumerate(kpis):
        x = M_L + i * (cw + 0.1)
        y = BODY_Y + 0.2
        add_kpi_card(s, x, y, cw, ch, big, label)

    # 중간 타이틀
    add_text(s, M_L, BODY_Y + 2.2, CONTENT_W, 0.4,
             '우리에게 부족한 단 한 가지',
             size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # 하단 큰 인용
    qy = BODY_Y + 2.7
    qh = 1.4
    add_rect(s, M_L, qy, CONTENT_W, qh, fill=ACCENT, line=None)
    add_text(s, M_L + 0.3, qy, CONTENT_W - 0.6, qh,
             '"부족한 것은 기술이 아니다.\n고객이 우리에게 닿는 길 — 디지털 창구 하나다."',
             size=18, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER, line_spacing=1.4)


# ========== CH3. ISS 2.0 위에 탑재 ==========

def slide_09_on_iss(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 03 · ON ISS 2.0',
               '별도 사이트를 만들지 않는다',
               'ISS 2.0 위에 얹어 진입장벽을 0으로', 9, TOTAL)

    # 좌: 기존 ISS 2.0 구조 / 우: 추가 레이어
    lw = CONTENT_W * 0.45
    rw = CONTENT_W - lw - 0.4

    add_text(s, M_L, BODY_Y + 0.1, lw, 0.3,
             '기존 — ISS 2.0 (Integrated Smart Ship)', size=11, bold=True, color=MUTED)

    layers = [
        ('① 데이터 수집', '엔진·주요기기 실시간 센서', BG_SOFT),
        ('② 분석·관제', '디지털 관제센터 + 부산 스마트케어', BG_SOFT),
        ('③ 리포트', '분기별 선주 제공 분석 리포트', BG_SOFT),
    ]
    ly = BODY_Y + 0.5
    lh = 0.85
    for i, (t, d, bg) in enumerate(layers):
        y = ly + i * (lh + 0.1)
        add_rect(s, M_L, y, lw, lh, fill=bg, line=LINE_SOFT)
        add_text(s, M_L + 0.18, y + 0.12, lw - 0.3, 0.3,
                 t, size=11, bold=True, color=INK)
        add_text(s, M_L + 0.18, y + 0.42, lw - 0.3, 0.4,
                 d, size=10, color=TEXT)

    # 화살표
    ax = M_L + lw + 0.05
    add_arrow(s, ax, BODY_Y + 1.7, ax + 0.3, BODY_Y + 1.7, color=ACCENT, weight=3)

    # 우측 — 새 레이어
    rx = M_L + lw + 0.4
    add_text(s, rx, BODY_Y + 0.1, rw, 0.3,
             '추가 — 24/7 무상 기술 지원 모듈', size=11, bold=True, color=ACCENT)

    new_layers = [
        ('④ 지원 요청', '1-click 영상 세션 시작 — Lv.1~3'),
        ('⑤ AI 정리', '세션 → Work Scope → PM 카드 자동 생성'),
        ('⑥ 이력 적층', 'CBM 데이터와 통합되어 자동 축적'),
    ]
    for i, (t, d) in enumerate(new_layers):
        y = ly + i * (lh + 0.1)
        add_rect(s, rx, y, rw, lh, fill=ACCENT, line=None)
        add_text(s, rx + 0.18, y + 0.12, rw - 0.3, 0.3,
                 t, size=11, bold=True, color=BG_WHITE)
        add_text(s, rx + 0.18, y + 0.42, rw - 0.3, 0.4,
                 d, size=10, color=BG_WHITE)

    add_takeaway(s, 'ISS 2.0 사용 선사는 추가 가입 없이 — 켜기만 하면 된다')


def slide_10_why_iss(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 03 · ON ISS 2.0',
               '왜 ISS 2.0 위에 얹는가',
               '이미 있는 데이터·로그인·고객을 그대로 활용한다', 10, TOTAL)

    reasons = [
        ('01', '기존 사용 선사 즉시 흡수',
         'ISS 2.0을 쓰는 선사는 이미\n로그인·인증·계약이 정리되어 있다.\n새 가입 동의 없이 모듈만 켜면 된다.'),
        ('02', '호선·점검 데이터 이미 존재',
         '엔진·기기 센서 데이터, 점검 이력,\n부품 정보가 ISS 2.0에 축적되어 있다.\n빈 사이트를 채울 필요가 없다.'),
        ('03', '별도 가입·로그인 없음',
         '오아시스가 부담스러웠던 이유 1순위.\n선주는 평소 보던 화면에서\n"지원 요청" 버튼만 누르면 된다.'),
        ('04', 'CBM ↔ 지원 세션 자동 연결',
         '알람 발생 화면에서 그대로 세션을 열고,\n조치 내역이 다시 ISS 2.0 이력에 쌓인다.\n맥락이 끊기지 않는다.'),
    ]
    cols = 2
    cw = (CONTENT_W - 0.25) / cols
    ch = 1.85
    for i, (num, t, d) in enumerate(reasons):
        r = i // cols
        c = i % cols
        x = M_L + c * (cw + 0.25)
        y = BODY_Y + 0.15 + r * (ch + 0.2)
        add_card(s, x, y, cw, ch,
                 kicker=f'WHY {num}', title=t, body=d,
                 accent_bar=True, title_size=14, body_size=10)

    add_takeaway(s, '새 사이트를 짓지 않고 — ISS 2.0이라는 기존 입구를 그대로 쓴다')


def slide_11_iss_synergy(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 03 · ON ISS 2.0',
               'ISS 2.0 × 기술 지원의 시너지',
               '고장 알람 → 즉시 원격 세션 → 기록 자동 축적', 11, TOTAL)

    # 가로 5단계 타임라인
    steps = [
        ('1', '알람 감지', 'ISS 2.0 센서가\n이상치 검출'),
        ('2', '1-click 요청', '알람 화면에서\n바로 세션 시작'),
        ('3', '원격 세션', '영상·화면공유\nLv.1~3 대응'),
        ('4', '조치 기록', 'AI가 자동 요약\n+ 스크린샷 보존'),
        ('5', 'CBM 적층', 'ISS 2.0 이력에\n자동 저장'),
    ]
    n = len(steps)
    sw = (CONTENT_W - 0.4 * (n - 1)) / n
    sh = 2.3
    sy = BODY_Y + 0.4
    for i, (num, t, d) in enumerate(steps):
        x = M_L + i * (sw + 0.4)
        # 본체
        is_last = i == n - 1
        fill = ACCENT if i in (0, n - 1) else BG_WHITE
        text_color = BG_WHITE if i in (0, n - 1) else INK
        body_color = BG_WHITE if i in (0, n - 1) else TEXT
        add_rect(s, x, sy, sw, sh, fill=fill, line=LINE_SOFT)
        # 번호 원
        add_text(s, x + 0.15, sy + 0.12, 0.4, 0.4,
                 num, size=20, bold=True, color=text_color)
        add_text(s, x + 0.15, sy + 0.7, sw - 0.3, 0.45,
                 t, size=13, bold=True, color=text_color)
        add_text(s, x + 0.15, sy + 1.2, sw - 0.3, 0.95,
                 d, size=10, color=body_color, line_spacing=1.4)
        # 화살표 (마지막 제외)
        if not is_last:
            ax1 = x + sw + 0.05
            ax2 = x + sw + 0.35
            add_arrow(s, ax1, sy + sh / 2, ax2, sy + sh / 2,
                     color=ACCENT_3, weight=2)

    add_takeaway(s, '선주는 이미 보고 있는 화면에서 도움을 받는다 — 새 사이트로 옮기지 않는다')


# ========== CH4. 24/7 무상 기술 지원 ==========

def slide_12_global_relay(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 04 · 24/7 SUPPORT',
               '해외지사를 시간대로 묶는다',
               '부산 ↔ 싱가포르 ↔ 유럽, 24시간 릴레이', 12, TOTAL)

    # 3개 거점 카드
    bases = [
        ('KST', '부산 본사', '00:00 — 08:00 KST\n야간/새벽 당직', '아시아 동부'),
        ('SGT', '싱가포르 지사', '08:00 — 16:00 KST\n주간 당직', '동남아·중동·인도'),
        ('CET', '유럽 지사', '16:00 — 24:00 KST\n저녁 당직', '유럽·아프리카·미주'),
    ]
    cw = (CONTENT_W - 0.4) / 3
    ch = 2.3
    cy = BODY_Y + 0.2
    for i, (tz, name, time, region) in enumerate(bases):
        x = M_L + i * (cw + 0.2)
        add_rect(s, x, cy, cw, ch, fill=BG_WHITE, line=LINE_SOFT)
        add_rect(s, x, cy, cw, 0.45, fill=ACCENT, line=None)
        add_text(s, x + 0.15, cy, cw - 0.3, 0.45,
                 tz, size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.18, cy + 0.6, cw - 0.36, 0.4,
                 name, size=15, bold=True, color=INK)
        add_text(s, x + 0.18, cy + 1.05, cw - 0.36, 0.65,
                 time, size=11, bold=True, color=ACCENT, line_spacing=1.3)
        add_text(s, x + 0.18, cy + 1.7, cw - 0.36, 0.5,
                 f'커버 지역  ·  {region}', size=10, color=TEXT)

    # SLA 박스
    sla_y = BODY_Y + 2.7
    add_rect(s, M_L, sla_y, CONTENT_W, 0.85, fill=BG_BOX, line=None)
    add_rect(s, M_L, sla_y, 0.05, 0.85, fill=ACCENT, line=None)
    add_text(s, M_L + 0.2, sla_y + 0.1, CONTENT_W - 0.3, 0.3,
             '응답 SLA (Service Level Agreement)', size=10, bold=True, color=ACCENT)
    add_text(s, M_L + 0.2, sla_y + 0.4, CONTENT_W * 0.5, 0.4,
             '일반 요청  ·  4시간 이내 세션 시작',
             size=14, bold=True, color=INK)
    add_text(s, M_L + CONTENT_W * 0.55, sla_y + 0.4, CONTENT_W * 0.45, 0.4,
             '긴급 요청  ·  1시간 이내 세션 시작',
             size=14, bold=True, color=WARN)

    add_takeaway(s, '선박이 어디 있든, 한국 시간 무관하게 — 4시간 또는 1시간')


def slide_13_three_levels(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 04 · 24/7 SUPPORT',
               '3단계 대응 체계',
               '내일 시작, 12개월 후 AR', 13, TOTAL)

    levels = [
        {
            'badge': 'LEVEL 1',
            'title': '전화 · 영상통화',
            'when': '즉시 도입',
            'cost': '추가 투자  ·  ₩0',
            'desc': '· 스마트폰 + Teams/WhatsApp\n· 선원이 촬영 → 엔지니어 가이드\n· 간단 문의·스위치 조작·알람 확인',
            'check': ['녹화', 'AR 마킹', '핸즈프리'],
            'has': [False, False, False],
            'is_main': False,
        },
        {
            'badge': 'LEVEL 2',
            'title': '화상 · 화면공유',
            'when': '3~6개월',
            'cost': '플랫폼 라이선스  ·  낮음',
            'desc': '· VSight / Teams Remote Assist\n· 화면 캡처·녹화·파일 전송\n· 중간 난이도 트러블슈팅',
            'check': ['녹화', 'AR 마킹', '핸즈프리'],
            'has': [True, False, False],
            'is_main': True,
        },
        {
            'badge': 'LEVEL 3',
            'title': 'AR 원격 가이드',
            'when': '12개월~',
            'cost': 'AR 기기 + 플랫폼  ·  중간',
            'desc': '· RealWear · HoloLens · 태블릿 AR\n· 시야에 화살표·텍스트 오버레이\n· 복잡한 수리·부품 교체 가이드',
            'check': ['녹화', 'AR 마킹', '핸즈프리'],
            'has': [True, True, True],
            'is_main': False,
        },
    ]

    cw = (CONTENT_W - 0.4) / 3
    ch = 4.0
    cy = BODY_Y + 0.15
    for i, lv in enumerate(levels):
        x = M_L + i * (cw + 0.2)
        is_main = lv['is_main']
        fill = ACCENT if is_main else BG_WHITE
        line = None if is_main else LINE_SOFT
        add_rect(s, x, cy, cw, ch, fill=fill, line=line)

        title_color = BG_WHITE if is_main else INK
        meta_color = BG_WHITE if is_main else MUTED
        body_color = BG_WHITE if is_main else TEXT
        accent_color = BG_WHITE if is_main else ACCENT

        # Badge
        add_text(s, x + 0.18, cy + 0.18, cw - 0.36, 0.3,
                 lv['badge'], size=9, bold=True, color=accent_color)
        add_text(s, x + 0.18, cy + 0.5, cw - 0.36, 0.4,
                 lv['title'], size=15, bold=True, color=title_color)
        add_text(s, x + 0.18, cy + 0.95, cw - 0.36, 0.3,
                 lv['when'], size=11, bold=True, color=accent_color)
        add_text(s, x + 0.18, cy + 1.25, cw - 0.36, 0.3,
                 lv['cost'], size=10, color=meta_color)

        add_line(s, x + 0.18, cy + 1.65, cw - 0.36, 0.005,
                color=BG_WHITE if is_main else LINE_SOFT)
        add_text(s, x + 0.18, cy + 1.75, cw - 0.36, 1.4,
                 lv['desc'], size=10, color=body_color, line_spacing=1.5)

        # Check list
        check_y = cy + ch - 1.0
        for j, (cap, has) in enumerate(zip(lv['check'], lv['has'])):
            mark = '●' if has else '○'
            color_mark = (BG_WHITE if is_main else ACCENT) if has else (
                BG_WHITE if is_main else LIGHT_MUTED)
            add_text(s, x + 0.18, check_y + j * 0.27, 0.25, 0.27,
                     mark, size=11, bold=True, color=color_mark, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, x + 0.42, check_y + j * 0.27, cw - 0.5, 0.27,
                     cap, size=10, color=body_color, anchor=MSO_ANCHOR.MIDDLE)

    add_takeaway(s, 'Lv.1은 추가 투자 없이 내일 시작 — Lv.2·3은 단계적 확장의 길')


def slide_14_wartsila_proof(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 04 · 24/7 SUPPORT',
               'Wärtsilä는 95%를 원격으로 푼다',
               '우리도 스마트폰 하나면 시작한다', 14, TOTAL)

    # 좌 — 큰 95% 박스
    lw = CONTENT_W * 0.42
    add_rect(s, M_L, BODY_Y + 0.2, lw, 3.4, fill=ACCENT, line=None)
    add_text(s, M_L + 0.3, BODY_Y + 0.45, lw - 0.6, 0.4,
             'WÄRTSILÄ SMART SUPPORT CENTRE',
             size=10, bold=True, color=BG_WHITE)
    add_text(s, M_L + 0.3, BODY_Y + 0.95, lw - 0.6, 1.5,
             '95%',
             size=120, bold=True, color=BG_WHITE,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_text(s, M_L + 0.3, BODY_Y + 2.5, lw - 0.6, 0.5,
             '오류를 원격으로 해결',
             size=16, bold=True, color=BG_WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, M_L + 0.3, BODY_Y + 3.0, lw - 0.6, 0.4,
             '"다음 항구에서 엔지니어를 부를 필요가 없다"',
             size=10, color=BG_WHITE,
             align=PP_ALIGN.CENTER)

    # 우 — 보조 사례 카드
    rx = M_L + lw + 0.3
    rw = CONTENT_W - lw - 0.3
    cases = [
        ('AR 정비 속도', '47%', 'AR 가이드 시 수리 속도 향상\n(종이 매뉴얼 대비)'),
        ('작업 오류율', '13%', 'AR 사용 시 오류율\n(종이 매뉴얼 53% 대비)'),
        ('미·캐 해군', '실전', 'ARMS · MIRRAS\n실전 배치 완료'),
        ('Kongsberg', 'PKI', 'Remote Diagnostics\n첫 시도 해결 원칙'),
    ]
    cw = (rw - 0.15) / 2
    ch = 1.6
    for i, (label, big, desc) in enumerate(cases):
        r = i // 2
        c = i % 2
        x = rx + c * (cw + 0.15)
        y = BODY_Y + 0.2 + r * (ch + 0.15)
        add_rect(s, x, y, cw, ch, fill=BG_WHITE, line=LINE_SOFT)
        add_text(s, x + 0.15, y + 0.12, cw - 0.3, 0.3,
                 label, size=9, bold=True, color=ACCENT)
        add_text(s, x + 0.15, y + 0.42, cw - 0.3, 0.55,
                 big, size=24, bold=True, color=INK)
        add_text(s, x + 0.15, y + 1.0, cw - 0.3, 0.55,
                 desc, size=10, color=TEXT, line_spacing=1.35)

    add_takeaway(s, '미래 기술이 아니다 — 이미 작동 중인 현재 기술이다')


def slide_15_why_free(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 04 · 24/7 SUPPORT',
               '무상으로 푸는 진짜 이유',
               '무료가 아니라 진입 미끼다', 15, TOTAL)

    # 펀넬 4단
    funnel = [
        ('STEP 1', '무상 24/7 접근',
         '진입장벽 없음 → 선주가 일단 들어온다',
         1.0, ACCENT_4),
        ('STEP 2', '신뢰 형성',
         '실제 응답·해결 경험이 누적된다',
         0.85, ACCENT_3),
        ('STEP 3', '데이터 축적',
         '세션·이력·CBM이 우리에게 쌓인다',
         0.7, ACCENT_2),
        ('STEP 4', '유상 전환',
         'LTSA Tier 1 → 2 → 3로 자연 이동',
         0.55, ACCENT),
    ]
    fy = BODY_Y + 0.3
    fh = 0.7
    fx_center = M_L + CONTENT_W / 2
    for i, (step, title, desc, ratio, color) in enumerate(funnel):
        bw = CONTENT_W * ratio
        x = fx_center - bw / 2
        y = fy + i * (fh + 0.12)
        add_rect(s, x, y, bw, fh, fill=color, line=None)
        add_text(s, x + 0.2, y, bw * 0.18, fh,
                 step, size=10, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + bw * 0.2, y, bw * 0.32, fh,
                 title, size=14, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + bw * 0.52, y, bw * 0.46, fh,
                 desc, size=10, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_takeaway(s, '오아시스의 가격 비교가 아니다 — 신뢰의 거리를 좁히는 것이 핵심이다')


# ========== CH5. PM 인수인계 ==========

def slide_16_handoff_flow(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 05 · PM HANDOFF',
               '지원 종료가 끝이 아니다',
               '세션 → 카드 → PM 자동 배정까지', 16, TOTAL)

    steps = [
        ('1', '세션 종료', '엔지니어가\n원격 세션 마감'),
        ('2', 'AI 요약', '대화·영상\n자동 분석 (3초)'),
        ('3', 'Work Scope', '일정·인원\n견적 자동 추출'),
        ('4', '건별 카드', '카드 생성 +\n타임라인 기록'),
        ('5', 'PM 배정', '담당 PM\n자동 매칭'),
        ('6', '알림 발송', '메일·웹·\n메신저 동시'),
    ]
    n = 6
    sw = (CONTENT_W - 0.2 * (n - 1)) / n
    sh = 2.5
    sy = BODY_Y + 0.4
    for i, (num, t, d) in enumerate(steps):
        x = M_L + i * (sw + 0.2)
        is_terminal = i in (0, n - 1)
        fill = ACCENT if is_terminal else BG_WHITE
        line = None if is_terminal else LINE_SOFT
        add_rect(s, x, sy, sw, sh, fill=fill, line=line)
        text_c = BG_WHITE if is_terminal else INK
        body_c = BG_WHITE if is_terminal else TEXT
        kicker_c = BG_WHITE if is_terminal else ACCENT

        # 번호
        add_text(s, x + 0.15, sy + 0.15, sw - 0.3, 0.6,
                 num, size=32, bold=True, color=kicker_c)
        add_text(s, x + 0.15, sy + 0.95, sw - 0.3, 0.45,
                 t, size=12, bold=True, color=text_c)
        add_text(s, x + 0.15, sy + 1.45, sw - 0.3, 1.0,
                 d, size=9, color=body_c, line_spacing=1.4)

        if i < n - 1:
            add_arrow(s, x + sw + 0.02, sy + sh / 2,
                     x + sw + 0.18, sy + sh / 2,
                     color=ACCENT_3, weight=2)

    add_takeaway(s, 'PM은 정리하지 않는다 — 정리된 카드를 받아 클로징만 한다')


def slide_17_ai_five(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 05 · PM HANDOFF',
               'AI가 정리하는 5가지',
               'PM이 메일을 다시 읽지 않는다', 17, TOTAL)

    items = [
        ('01', '메일·세션 자동 분석',
         '수신 메일과 영상 세션 텍스트를\n3초 이내 파싱하여 핵심 항목을 추출한다.',
         '응답 시간  ·  3초'),
        ('02', '건별 카드 자동 생성',
         '분석 결과로 작업 카드를 만들고,\n타임라인에 자동 등록한다.',
         '입력 작업  ·  0건'),
        ('03', '협력업체·러프 견적 매칭',
         '과거 이력 기반으로 적합 업체를 추천하고,\n초안 견적가를 산출한다.',
         '견적 회신  ·  당일'),
        ('04', '수기발주서 자동 드래프트',
         '긴급공사 + PO 미확인 상황을 감지해\n수기발주서 초안과 알림을 발송한다.',
         '하도급법  ·  자동 대응'),
        ('05', 'PM 변경 시 인수인계 문서',
         '담당 건 전체를 AI가 요약 문서화하여\n새 PM이 30분 내 파악하도록 한다.',
         '인수인계  ·  30분'),
    ]
    # 1열에 3개, 2열에 2개 (3+2 그리드 변형)
    # 더 깔끔하게: 5개를 1행 5열로 배치하면 좁아지므로 2행 (3+2)
    cw = (CONTENT_W - 0.2 * 2) / 3
    ch = (BODY_H - 0.5 - 0.15) / 2
    positions = [(0, 0), (0, 1), (0, 2), (1, 0), (1, 1)]
    for (r, c), (num, t, d, meta) in zip(positions, items):
        x = M_L + c * (cw + 0.2)
        y = BODY_Y + 0.15 + r * (ch + 0.15)
        # 5번째는 가운데 정렬
        if r == 1 and c == 1:
            x = M_L + cw + 0.2 + (cw + 0.2) / 2
        add_rect(s, x, y, cw, ch, fill=BG_WHITE, line=LINE_SOFT)
        add_rect(s, x, y, 0.06, ch, fill=ACCENT, line=None)
        add_text(s, x + 0.18, y + 0.12, cw - 0.3, 0.25,
                 f'AI · {num}', size=9, bold=True, color=ACCENT)
        add_text(s, x + 0.18, y + 0.42, cw - 0.3, 0.4,
                 t, size=13, bold=True, color=INK)
        add_text(s, x + 0.18, y + 0.92, cw - 0.3, 0.95,
                 d, size=10, color=TEXT, line_spacing=1.4)
        # 메타 라인
        meta_y = y + ch - 0.45
        add_line(s, x + 0.18, meta_y, cw - 0.36, 0.005, color=LINE_SOFT)
        add_text(s, x + 0.18, meta_y + 0.07, cw - 0.3, 0.3,
                 meta, size=10, bold=True, color=ACCENT)

    # 5번째 카드를 정 가운데로 보정 — 위 코드에서 c==1일 때 다른 x 적용했으니 별도 처리 불필요
    # (단, 첫 행에 3개가 들어가므로 두 번째 행은 자연스럽게 좌측부터 채워짐)


def slide_18_before_after(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 05 · PM HANDOFF',
               'Before · After',
               '메일 포워딩 1~2주 → AI 요약 30분', 18, TOTAL)

    rows = [
        ('인수인계', '메일 포워딩 → 사람이 읽고 정리 (1~2주)', 'AI 요약 카드 (30분)'),
        ('견적 회신', '담당자 메모리·과거 메일 검색 (2~3일)', 'AI 러프 견적 + 즉시 검토 (당일)'),
        ('진행 추적', '전화/메일로 PM에게 문의', '웹에서 단계별 실시간 표시'),
        ('노하우 보존', 'PM 머리·개인 메일함', '세션 녹화·카드·DB로 영구 축적'),
        ('하도급법 리스크', 'PM 기억에 의존, 누락 시 수기발주서 미발행', '시스템이 감지·드래프트·알림 자동 처리'),
    ]
    tx = M_L
    ty = BODY_Y + 0.2
    th = 0.45
    col_label = CONTENT_W * 0.18
    col_before = CONTENT_W * 0.45
    col_after = CONTENT_W * 0.37

    add_rect(s, tx, ty, col_label, th, fill=DARK_BG, line=None)
    add_rect(s, tx + col_label, ty, col_before, th, fill=DARK_BG, line=None)
    add_rect(s, tx + col_label + col_before, ty, col_after, th, fill=ACCENT, line=None)
    add_text(s, tx + 0.15, ty, col_label - 0.3, th,
             '항목', size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + col_label + 0.15, ty, col_before - 0.3, th,
             'BEFORE', size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + col_label + col_before + 0.15, ty, col_after - 0.3, th,
             'AFTER', size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)

    cur_y = ty + th
    rh = 0.62
    for i, (label, before, after) in enumerate(rows):
        bg = BG_SOFT if i % 2 == 0 else BG_WHITE
        add_rect(s, tx, cur_y, CONTENT_W, rh, fill=bg, line=None)
        add_line(s, tx, cur_y, CONTENT_W, 0.005, color=LINE_SOFT)
        add_text(s, tx + 0.15, cur_y, col_label - 0.3, rh,
                 label, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + col_label + 0.15, cur_y, col_before - 0.3, rh,
                 before, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
        add_text(s, tx + col_label + col_before + 0.15, cur_y, col_after - 0.3, rh,
                 after, size=10, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
        cur_y += rh

    add_takeaway(s, 'PM 시간을 돌려준다 = 더 많은 건을 더 빠르게 클로징한다')


# ========== CH6. 유상 전환 — 구독형 메인 ==========

def slide_19_two_paths(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 06 · MONETIZATION',
               '무상이 신뢰가 되면 유상이 된다',
               '메인은 구독형 LTSA — 건별은 보조 진입로', 19, TOTAL)

    # 가운데 시작점
    cx = M_L + CONTENT_W / 2
    sy = BODY_Y + 0.35
    start_w = 3.5
    start_h = 0.7
    add_rect(s, cx - start_w / 2, sy, start_w, start_h, fill=DARK_BG, line=None)
    add_text(s, cx - start_w / 2, sy, start_w, start_h,
             '무상 24/7 사용자 풀',
             size=14, bold=True, color=BG_WHITE,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # 두 갈래 화살표
    branch_y_top = sy + start_h + 0.1
    branch_y_btm = sy + start_h + 0.7
    add_arrow(s, cx, branch_y_top, M_L + CONTENT_W * 0.25, branch_y_btm,
             color=ACCENT, weight=3)
    add_arrow(s, cx, branch_y_top, M_L + CONTENT_W * 0.75, branch_y_btm,
             color=MUTED, weight=2)

    # 좌측 — 구독형 (메인)
    bw = CONTENT_W * 0.42
    by = branch_y_btm + 0.1
    bh = 2.6
    bx = M_L
    add_rect(s, bx, by, bw, bh, fill=ACCENT, line=None)
    add_text(s, bx + 0.25, by + 0.18, bw - 0.5, 0.3,
             'PATH A — 메인', size=10, bold=True, color=BG_WHITE)
    add_text(s, bx + 0.25, by + 0.5, bw - 0.5, 0.5,
             '구독형 LTSA',
             size=22, bold=True, color=BG_WHITE)
    add_text(s, bx + 0.25, by + 1.05, bw - 0.5, 0.4,
             '라이프사이클 계약 3티어',
             size=12, color=BG_WHITE)
    add_text(s, bx + 0.25, by + 1.55, bw - 0.5, 1.0,
             '· 다수 선박 보유 선주\n· 정기 AS 패턴이 있는 고객\n· 비용 예측·우선 대응이 필요한 고객',
             size=11, color=BG_WHITE, line_spacing=1.45)

    # 우측 — 건별 (보조)
    rx = M_L + CONTENT_W - bw
    add_rect(s, rx, by, bw, bh, fill=BG_WHITE, line=LINE_SOFT)
    add_text(s, rx + 0.25, by + 0.18, bw - 0.5, 0.3,
             'PATH B — 보조', size=10, bold=True, color=MUTED)
    add_text(s, rx + 0.25, by + 0.5, bw - 0.5, 0.5,
             '건별 Pay-as-you-go',
             size=22, bold=True, color=INK)
    add_text(s, rx + 0.25, by + 1.05, bw - 0.5, 0.4,
             '비계약 고객의 진입로',
             size=12, color=MUTED)
    add_text(s, rx + 0.25, by + 1.55, bw - 0.5, 1.0,
             '· 단일 선박 보유 선주\n· 비정기·돌발 수요 고객\n· 향후 구독 전환 후보',
             size=11, color=TEXT, line_spacing=1.45)

    add_takeaway(s, '선주가 고르는 게 아니라, 사용 패턴이 자연스럽게 안내한다')


def slide_20_lts_tiers(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 06 · MONETIZATION',
               '구독형 — 라이프사이클 계약 3티어',
               'Wärtsilä Service Value Ladder Lv.4 진입', 20, TOTAL)

    tiers = [
        {
            'badge': 'TIER 1', 'name': 'Basic',
            'tag': '정기 점검 계약',
            'price': '진입형',
            'period': '계약 1~2년',
            'feats': '· 연 1~2회 헬스체크\n· 시각화 리포트\n· 웹 포털 접근\n· 점검 알림\n· 호선 이력 조회',
            'is_main': False,
        },
        {
            'badge': 'TIER 2', 'name': 'Standard',
            'tag': '종합 관리 계약',
            'price': '핵심 티어',
            'period': '계약 2~3년',
            'feats': '· Tier 1 전체\n· 무제한 기술 지원 (24h SLA)\n· 우선 대응 + 전담 PM\n· 부품 수명 관리\n· AI 분석 리포트',
            'is_main': True,
        },
        {
            'badge': 'TIER 3', 'name': 'Premium',
            'tag': '토탈 케어 계약',
            'price': '최상위',
            'period': '계약 3~5년',
            'feats': '· Tier 2 전체\n· 수리 비용 정액제 포함\n· 긴급 대응 4h SLA\n· 원격 AR 기술 지원\n· 정비 계획 수립 + 교육',
            'is_main': False,
        },
    ]
    cw = (CONTENT_W - 0.4) / 3
    ch = 4.0
    cy = BODY_Y + 0.15
    for i, t in enumerate(tiers):
        x = M_L + i * (cw + 0.2)
        is_main = t['is_main']
        fill = ACCENT if is_main else BG_WHITE
        line = None if is_main else LINE_SOFT
        add_rect(s, x, cy, cw, ch, fill=fill, line=line)

        ink_c = BG_WHITE if is_main else INK
        muted_c = BG_WHITE if is_main else MUTED
        accent_c = BG_WHITE if is_main else ACCENT
        body_c = BG_WHITE if is_main else TEXT

        # 추천 배지
        if is_main:
            badge_w = 1.5
            add_rect(s, x + cw - badge_w - 0.15, cy + 0.15, badge_w, 0.32,
                    fill=BG_WHITE, line=None)
            add_text(s, x + cw - badge_w - 0.15, cy + 0.15, badge_w, 0.32,
                     'RECOMMENDED', size=9, bold=True, color=ACCENT,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

        add_text(s, x + 0.2, cy + 0.18, cw - 0.4, 0.3,
                 t['badge'], size=10, bold=True, color=accent_c)
        add_text(s, x + 0.2, cy + 0.5, cw - 0.4, 0.55,
                 t['name'], size=24, bold=True, color=ink_c)
        add_text(s, x + 0.2, cy + 1.1, cw - 0.4, 0.3,
                 t['tag'], size=12, color=muted_c)
        add_line(s, x + 0.2, cy + 1.45, cw - 0.4, 0.005,
                color=BG_WHITE if is_main else LINE_SOFT)
        add_text(s, x + 0.2, cy + 1.55, cw - 0.4, 0.3,
                 t['price'], size=11, bold=True, color=accent_c)
        add_text(s, x + 0.2, cy + 1.85, cw - 0.4, 0.3,
                 t['period'], size=10, color=body_c)
        add_line(s, x + 0.2, cy + 2.2, cw - 0.4, 0.005,
                color=BG_WHITE if is_main else LINE_SOFT)
        add_text(s, x + 0.2, cy + 2.32, cw - 0.4, 1.55,
                 t['feats'], size=10, color=body_c, line_spacing=1.55)

    add_takeaway(s, 'Tier 1로 진입 → 데이터로 신뢰 → Tier 2 → 3로 자연 상승하는 구조')


def slide_21_pay_as_you_go(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 06 · MONETIZATION',
               '건별 — Pay-as-you-go',
               '비계약 고객도 그대로 우리 안에 둔다', 21, TOTAL)

    # 좌 — 건별 과금 항목
    lw = CONTENT_W * 0.5
    add_text(s, M_L, BODY_Y + 0.1, lw, 0.3,
             '건별 과금 항목', size=11, bold=True, color=ACCENT)
    items = [
        ('원격 세션', 'Lv.1~2 시간당', '시간당 정액'),
        ('현장 출장', '왕복+체류+공임', '실비 + 공임'),
        ('리포트 단건', '시각화 분석 리포트', '리포트당 정액'),
        ('기술 문서', '도면·매뉴얼 단건', '문서당 소액'),
    ]
    cols = 2
    cw_item = (lw - 0.15) / cols
    ch_item = 1.0
    for i, (t, d, price) in enumerate(items):
        r = i // cols
        c = i % cols
        x = M_L + c * (cw_item + 0.15)
        y = BODY_Y + 0.5 + r * (ch_item + 0.12)
        add_rect(s, x, y, cw_item, ch_item, fill=BG_WHITE, line=LINE_SOFT)
        add_text(s, x + 0.18, y + 0.12, cw_item - 0.3, 0.3,
                 t, size=12, bold=True, color=INK)
        add_text(s, x + 0.18, y + 0.42, cw_item - 0.3, 0.3,
                 d, size=10, color=TEXT)
        add_text(s, x + 0.18, y + 0.7, cw_item - 0.3, 0.25,
                 price, size=10, bold=True, color=ACCENT)

    # 우 — 전환 트리거
    rx = M_L + lw + 0.3
    rw = CONTENT_W - lw - 0.3
    add_text(s, rx, BODY_Y + 0.1, rw, 0.3,
             '구독 전환 트리거', size=11, bold=True, color=ACCENT)
    triggers = [
        ('연 N건 이상 사용', '건별 누적 ₩30M 초과 → Tier 1 제안'),
        ('다중 선박 운용', '관리 선박 5척 이상 → Tier 2 제안'),
        ('긴급 수리 빈도', '긴급 SLA 요청 잦음 → Tier 3 제안'),
        ('정기 AS 패턴', '연 헬스체크 정기화 → Tier 1 자동 권유'),
    ]
    ty = BODY_Y + 0.5
    th = 0.55
    for i, (t, d) in enumerate(triggers):
        y = ty + i * (th + 0.1)
        add_rect(s, rx, y, rw, th, fill=BG_BOX, line=None)
        add_rect(s, rx, y, 0.05, th, fill=ACCENT, line=None)
        add_text(s, rx + 0.2, y, rw * 0.4, th,
                 t, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + rw * 0.42, y, rw * 0.55, th,
                 d, size=10, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    add_takeaway(s, '건별 사용자는 떠나는 고객이 아니라 — 미래의 구독자다')


def slide_22_subscription_vs_payg(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 06 · MONETIZATION',
               '두 모델 비교',
               '구조적으로 구독형이 우월하다', 22, TOTAL)

    # 비교표
    tx = M_L
    ty = BODY_Y + 0.2
    th = 0.45
    col_label = CONTENT_W * 0.22
    col_sub = CONTENT_W * 0.39
    col_pay = CONTENT_W * 0.39

    add_rect(s, tx, ty, col_label, th, fill=DARK_BG, line=None)
    add_rect(s, tx + col_label, ty, col_sub, th, fill=ACCENT, line=None)
    add_rect(s, tx + col_label + col_sub, ty, col_pay, th, fill=DARK_BG, line=None)
    add_text(s, tx + 0.15, ty, col_label - 0.3, th,
             '평가 항목', size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + col_label + 0.15, ty, col_sub - 0.3, th,
             '구독형 LTSA', size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + col_label + col_sub + 0.15, ty, col_pay - 0.3, th,
             '건별 Pay-as-you-go', size=11, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ('매출 예측', '연간 계약금 확정 — 월/분기 안정', '고장 발생 시점에 의존'),
        ('우선 대응', 'SLA 보증 (4~24시간)', '없음 — 일반 큐'),
        ('데이터 축적', '연속 데이터 — 우리만 가진 자산', '단편적 — 건별 분리'),
        ('마진 구조', '패키지 프리미엄 가능', '매번 가격 비교 대상'),
        ('Lock-in 효과', '계약 기간 내 이탈 없음', '없음 — 다음 건 경쟁사로'),
    ]
    cur_y = ty + th
    rh = 0.5
    for i, (label, sub, pay) in enumerate(rows):
        bg = BG_SOFT if i % 2 == 0 else BG_WHITE
        add_rect(s, tx, cur_y, CONTENT_W, rh, fill=bg, line=None)
        add_line(s, tx, cur_y, CONTENT_W, 0.005, color=LINE_SOFT)
        add_text(s, tx + 0.15, cur_y, col_label - 0.3, rh,
                 label, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + col_label + 0.15, cur_y, col_sub - 0.3, rh,
                 sub, size=10, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + col_label + col_sub + 0.15, cur_y, col_pay - 0.3, rh,
                 pay, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
        cur_y += rh

    # 시뮬레이션 박스
    sy = cur_y + 0.15
    add_rect(s, M_L, sy, CONTENT_W, 0.6, fill=ACCENT, line=None)
    add_text(s, M_L + 0.3, sy, CONTENT_W - 0.6, 0.6,
             '시뮬레이션  ·  Tier 2 ₩20M × 50척 = 연 ₩10억 반복 매출 — 고장 유무와 무관하게 들어온다',
             size=13, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)


# ========== CH7. 기대효과 ==========

def slide_23_quantitative(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 07 · IMPACT',
               '정량 효과',
               '숫자로 증명되는 8가지', 23, TOTAL)

    kpis = [
        ('출장 비용', '▼50%', '연 ₩150M 절감 (100건 시뮬)'),
        ('견적 회신', '당일', '기존 2~3일 → 당일 (AI 자동 산출)'),
        ('PM 인수인계', '30분', '기존 1~2주 → 30분 (AI 요약)'),
        ('해외 대응', '수 시간', '기존 1~2주 → 수 시간 (원격)'),
        ('원격 해결률', '95%', 'Wärtsilä Smart Support Centre 실증'),
        ('수리 속도', '▲47%', 'AR 가이드 시 (종이 매뉴얼 대비)'),
        ('엔지니어 처리량', '×3~5', '같은 인원으로 더 많은 건 처리'),
        ('Tier 2 매출', '₩10억', '50척 가입 시 연간 반복 매출'),
    ]
    cols = 4
    cw = (CONTENT_W - 0.15 * (cols - 1)) / cols
    ch = (BODY_H - 0.3 - 0.2) / 2
    for i, (label, big, desc) in enumerate(kpis):
        r = i // cols
        c = i % cols
        x = M_L + c * (cw + 0.15)
        y = BODY_Y + 0.15 + r * (ch + 0.15)
        add_rect(s, x, y, cw, ch, fill=BG_WHITE, line=LINE_SOFT)
        add_rect(s, x, y, cw, 0.06, fill=ACCENT, line=None)
        add_text(s, x + 0.15, y + 0.18, cw - 0.3, 0.3,
                 label, size=10, bold=True, color=ACCENT)
        add_text(s, x + 0.15, y + 0.55, cw - 0.3, 0.85,
                 big, size=32, bold=True, color=INK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.15, y + ch - 0.7, cw - 0.3, 0.6,
                 desc, size=9, color=TEXT, line_spacing=1.4)


def slide_24_qualitative(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 07 · IMPACT',
               '정성 효과',
               '고객 · 회사 · 그룹 — 3중 효과', 24, TOTAL)

    cols_data = [
        ('고객', 'CUSTOMER', '선주의 일상이 바뀐다',
         '· 24/7 어디서든 접근\n· 시각화 리포트로 빠른 의사결정\n· 비용을 미리 예측 가능\n· 우리가 먼저 알려주는 알림\n· 매뉴얼·도면 셀프서비스'),
        ('회사', 'COMPANY', '우리의 사업 구조가 바뀐다',
         '· 안정적 반복 매출\n· 데이터 기반 선제 영업\n· PM 시간 회수 → 처리량 확대\n· 고객 이탈 방지 (Lock-in)\n· 데이터 자산 축적'),
        ('그룹', 'GROUP', 'HD현대 AX 전략과 정합한다',
         '· Hi-4S · ISS 2.0과 자연 연동\n· 모회사 디지털 전환 가속\n· 9,890척 자산 활용 극대화\n· 글로벌 디지털 표준 진입\n· OEM 정통성 강화'),
    ]
    cw = (CONTENT_W - 0.4) / 3
    ch = 4.0
    cy = BODY_Y + 0.15
    for i, (name_kr, name_en, sub, body) in enumerate(cols_data):
        x = M_L + i * (cw + 0.2)
        is_main = i == 1  # 회사 강조
        fill = ACCENT if is_main else BG_WHITE
        line = None if is_main else LINE_SOFT
        add_rect(s, x, cy, cw, ch, fill=fill, line=line)
        ink_c = BG_WHITE if is_main else INK
        accent_c = BG_WHITE if is_main else ACCENT
        body_c = BG_WHITE if is_main else TEXT
        muted_c = BG_WHITE if is_main else MUTED

        add_text(s, x + 0.25, cy + 0.25, cw - 0.5, 0.3,
                 name_en, size=10, bold=True, color=accent_c)
        add_text(s, x + 0.25, cy + 0.6, cw - 0.5, 0.55,
                 name_kr, size=22, bold=True, color=ink_c)
        add_text(s, x + 0.25, cy + 1.2, cw - 0.5, 0.4,
                 sub, size=11, color=muted_c)
        add_line(s, x + 0.25, cy + 1.7, cw - 0.5, 0.005,
                color=BG_WHITE if is_main else LINE_SOFT)
        add_text(s, x + 0.25, cy + 1.85, cw - 0.5, 2.0,
                 body, size=11, color=body_c, line_spacing=1.6)

    add_takeaway(s, '한 번 만든 디지털 창구가 — 고객 · 회사 · 그룹 세 방향으로 동시에 작동한다')


def slide_25_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CHAPTER 07 · IMPACT',
               '로드맵',
               'Phase 1은 내일 시작할 수 있다', 25, TOTAL)

    phases = [
        ('PHASE 1', '1~3개월',
         'Lv.1 원격 지원 +\nAI PM 자동화 +\n시각화 리포트 템플릿',
         '추가 투자 ₩0', ACCENT),
        ('PHASE 2', '3~6개월',
         'ISS 2.0 위 포털 MVP +\nLv.2 화상 플랫폼 +\nTier 1 파일럿',
         '플랫폼 라이선스', ACCENT_2),
        ('PHASE 3', '6~12개월',
         '대시보드 + 문서 라이브러리 +\n알림 시스템 +\nTier 1 → Tier 2 확대',
         '중간 규모 투자', ACCENT_3),
        ('PHASE 4', '12개월~',
         'Lv.3 AR 가이드 +\nTier 2 → Tier 3 +\n그룹 Hi-4S 연동',
         '그룹 협의', ACCENT_4),
    ]
    n = 4
    sw = (CONTENT_W - 0.2 * (n - 1)) / n
    sh = 3.6
    sy = BODY_Y + 0.2
    for i, (badge, period, body, cost, color) in enumerate(phases):
        x = M_L + i * (sw + 0.2)
        add_rect(s, x, sy, sw, 0.7, fill=color, line=None)
        add_text(s, x + 0.15, sy, sw - 0.3, 0.35,
                 badge, size=11, bold=True, color=BG_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.15, sy + 0.32, sw - 0.3, 0.35,
                 period, size=12, bold=True, color=BG_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)

        add_rect(s, x, sy + 0.7, sw, sh - 0.7, fill=BG_WHITE, line=LINE_SOFT)
        add_text(s, x + 0.18, sy + 0.85, sw - 0.36, 1.7,
                 body, size=11, color=INK, line_spacing=1.55)
        add_line(s, x + 0.18, sy + sh - 0.7, sw - 0.36, 0.005, color=LINE_SOFT)
        add_text(s, x + 0.18, sy + sh - 0.55, sw - 0.36, 0.3,
                 '투자 규모', size=9, bold=True, color=ACCENT)
        add_text(s, x + 0.18, sy + sh - 0.25, sw - 0.36, 0.25,
                 cost, size=11, bold=True, color=INK)

    add_takeaway(s, 'Phase 1은 추가 투자 없이 — 의사결정만 있으면 시작된다')


# ========== 마무리 ==========

def slide_26_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'CLOSING', '한 문장으로',
               '우리는 오아시스를 다시 짓는 것이 아니다', 26, TOTAL)

    # 큰 인용 카드 (단정조)
    qy = BODY_Y + 0.2
    qh = 3.0
    add_rect(s, M_L, qy, CONTENT_W, qh, fill=ACCENT, line=None)
    add_text(s, M_L + 0.5, qy + 0.4, CONTENT_W - 1.0, qh - 0.8,
             '오아시스는 중개를 팔다가 멈췄다.\n\n'
             '우리는 24/7 무상 기술 지원으로 신뢰를 만든다.\n'
             '그 신뢰는 라이프사이클 계약이 된다.\n'
             '그 계약은 우리만의 데이터로 쌓인다.\n\n'
             '이것이 OEM AS사가 가야 할 길이다.',
             size=20, bold=True, color=BG_WHITE,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, line_spacing=1.55)

    # 하단 미니 KPI 3종
    by = qy + qh + 0.25
    bh = 0.85
    bw = (CONTENT_W - 0.3) / 3
    minis = [
        ('9,890척', 'AS 네트워크가 곧 진입로다'),
        ('24/7', '시간이 곧 신뢰다'),
        ('ISS 2.0', '데이터가 곧 자산이다'),
    ]
    for i, (big, label) in enumerate(minis):
        x = M_L + i * (bw + 0.15)
        add_rect(s, x, by, bw, bh, fill=BG_WHITE, line=LINE_SOFT)
        add_text(s, x + 0.2, by, bw * 0.4, bh,
                 big, size=18, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + bw * 0.42, by, bw * 0.55, bh,
                 label, size=11, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)


def slide_27_qa(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_bg(s, BG_WHITE)
    add_header(s, 'Q&A', '질문 환영합니다', 'Thank You', 27, TOTAL)

    # 좌 — 발표자 카드 (높이 확장: 3.0 → 4.05)
    lw = CONTENT_W * 0.45
    lh = 4.05
    add_rect(s, M_L, BODY_Y + 0.15, lw, lh, fill=ACCENT, line=None)
    add_text(s, M_L + 0.3, BODY_Y + 0.4, lw - 0.6, 0.3,
             'PRESENTER', size=10, bold=True, color=BG_WHITE)
    add_text(s, M_L + 0.3, BODY_Y + 0.75, lw - 0.6, 0.5,
             'HD현대마린솔루션테크',
             size=18, bold=True, color=BG_WHITE)
    add_text(s, M_L + 0.3, BODY_Y + 1.3, lw - 0.6, 0.4,
             'AS 사업본부', size=13, color=BG_WHITE)
    add_line(s, M_L + 0.3, BODY_Y + 1.85, lw - 0.6, 0.005, color=BG_WHITE)

    add_text(s, M_L + 0.3, BODY_Y + 2.0, lw - 0.6, 0.3,
             '문의', size=10, bold=True, color=BG_WHITE)
    add_text(s, M_L + 0.3, BODY_Y + 2.3, lw - 0.6, 0.3,
             '내선 / 이메일 / 사내 메신저',
             size=11, color=BG_WHITE)

    add_text(s, M_L + 0.3, BODY_Y + 2.75, lw - 0.6, 0.3,
             '후속 액션', size=10, bold=True, color=BG_WHITE)
    add_text(s, M_L + 0.3, BODY_Y + 3.05, lw - 0.6, 1.0,
             '· Phase 1 의사결정 회의 (1주 내)\n· Lv.1 원격 지원 파일럿 셋업\n· Tier 1 후보 선주 3곳 선정',
             size=11, color=BG_WHITE, line_spacing=1.5)

    # 우 — 참고 자료
    rx = M_L + lw + 0.3
    rw = CONTENT_W - lw - 0.3
    add_text(s, rx, BODY_Y + 0.15, rw, 0.3,
             '참고 자료', size=11, bold=True, color=ACCENT)
    refs = [
        ('AM 벤치마킹', '선도기업 AM 전략 벤치마킹 보고서'),
        ('원격 지원', '원격 기술 지원 — 24/7 전문가 대기 체계 제안'),
        ('웹사이트 구축', '24/7 무상 기술지원 웹사이트 구축 제안'),
        ('LTSA', '라이프사이클 계약 도입 제안'),
        ('웹포털 설계', '웹 포털 상세 설계 + 격차 해소 방안'),
    ]
    ty = BODY_Y + 0.55
    th = 0.5
    for i, (label, desc) in enumerate(refs):
        y = ty + i * (th + 0.08)
        add_rect(s, rx, y, rw, th, fill=BG_BOX, line=None)
        add_rect(s, rx, y, 0.05, th, fill=ACCENT, line=None)
        add_text(s, rx + 0.2, y, rw * 0.3, th,
                 label, size=10, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + rw * 0.32, y, rw * 0.65, th,
                 desc, size=10, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    # 우 하단 — 클로징 한 줄 박스
    cy = BODY_Y + 3.6
    add_rect(s, rx, cy, rw, 0.55, fill=DARK_BG, line=None)
    add_text(s, rx + 0.2, cy, rw - 0.3, 0.55,
             '시작은 의사결정 한 번으로 충분합니다',
             size=12, bold=True, color=BG_WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ========== 메인 ==========

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    builders = [
        slide_01_cover, slide_02_agenda,
        slide_03_oasis_what, slide_04_oasis_why_failed, slide_05_oasis_lesson,
        slide_06_benchmark, slide_07_we_are_different, slide_08_gap_essence,
        slide_09_on_iss, slide_10_why_iss, slide_11_iss_synergy,
        slide_12_global_relay, slide_13_three_levels,
        slide_14_wartsila_proof, slide_15_why_free,
        slide_16_handoff_flow, slide_17_ai_five, slide_18_before_after,
        slide_19_two_paths, slide_20_lts_tiers,
        slide_21_pay_as_you_go, slide_22_subscription_vs_payg,
        slide_23_quantitative, slide_24_qualitative, slide_25_roadmap,
        slide_26_closing, slide_27_qa,
    ]
    assert len(builders) == TOTAL, f'슬라이드 빌더 개수 불일치: {len(builders)} vs {TOTAL}'

    for fn in builders:
        fn(prs)

    # 폰트 검증
    bad = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name and run.font.name != FONT_NAME:
                        bad += 1
    if bad:
        print(f'[WARN] 비-맑은고딕 run: {bad}')
    else:
        print('[OK] 모든 텍스트 run이 맑은 고딕으로 설정되었습니다.')

    out = '/Users/sinchaeyeon/Desktop/채린/24_7_무상_기술지원_웹사이트_제안.pptx'
    prs.save(out)
    print(f'[OK] 저장 완료: {out}')
    print(f'[OK] 슬라이드 수: {len(prs.slides)}')


if __name__ == '__main__':
    main()
