#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colors ──────────────────────────────────────────────────────────
C_BLACK    = RGBColor(0x0A, 0x0A, 0x0A)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_CORAL    = RGBColor(0xE5, 0x3A, 0x27)
C_CHARCOAL = RGBColor(0x37, 0x41, 0x51)
C_STEEL    = RGBColor(0x6B, 0x7B, 0x8D)
C_HAIRLINE = RGBColor(0xE0, 0xE2, 0xE6)
C_SURFACE  = RGBColor(0xF5, 0xF6, 0xF7)
C_MIDGRAY  = RGBColor(0x9A, 0x9A, 0x9A)

FONT = '맑은 고딕'
I    = Inches


# ── Helpers ─────────────────────────────────────────────────────────
def make_prs():
    prs = Presentation()
    prs.slide_width  = I(13.33)
    prs.slide_height = I(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def no_line(shp):
    spPr = shp._element.spPr
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    ln.clear()
    etree.SubElement(ln, qn('a:noFill'))

def rect(slide, l, t, w, h, fill_c=None, line_c=None, line_pt=0.75):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill_c:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_c
    else:
        shp.fill.background()
    if line_c:
        shp.line.color.rgb = line_c
        shp.line.width = Pt(line_pt)
    else:
        no_line(shp)
    return shp

def add_tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    return box.text_frame

def run_style(run, size, bold=False, color=C_BLACK):
    run.font.name  = FONT
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color

def first_para(tf, text, size, bold=False, color=C_BLACK, align=PP_ALIGN.LEFT,
               space_before=0, line_spacing=None):
    p = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    r = p.add_run()
    r.text = text
    run_style(r, size, bold, color)
    return p

def add_para(tf, text, size, bold=False, color=C_BLACK, align=PP_ALIGN.LEFT,
             space_before=0, line_spacing=None):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    r = p.add_run()
    r.text = text
    run_style(r, size, bold, color)
    return p


# ── Slide builders ──────────────────────────────────────────────────
W_FULL = I(13.33)
H_FULL = I(7.5)

# Left accent bar (coral strip, constant position across content slides)
ACCENT_L = I(0)
ACCENT_W = I(0.07)

# Chapter tag (top-left), slide title (below), divider line (below title)
TAG_L   = I(0.38)
TAG_T   = I(0.20)
TITLE_L = I(0.38)
TITLE_T = I(0.52)
DIV_T   = I(1.32)
DIV_H   = I(0.022)

# Content start
CONTENT_L = I(0.38)
CONTENT_W = I(12.8)


def slide_cover(prs, title_top, title_main, dept="전략기획팀", year="2026"):
    s = blank(prs)
    set_bg(s, C_BLACK)
    # Top thin coral bar
    rect(s, I(0), I(0), W_FULL, I(0.05), fill_c=C_CORAL)
    # Left coral bar
    rect(s, I(0), I(0), I(0.07), H_FULL, fill_c=C_CORAL)

    # Subtitle (small, top)
    f = add_tb(s, I(0.5), I(1.8), I(11.5), I(0.5))
    first_para(f, title_top, 16, False, C_MIDGRAY)

    # Main title
    f2 = add_tb(s, I(0.5), I(2.3), I(11.5), I(1.7))
    first_para(f2, title_main, 46, True, C_WHITE, line_spacing=58)

    # Coral accent line
    rect(s, I(0.5), I(4.35), I(2.2), I(0.04), fill_c=C_CORAL)

    # Date / dept
    f3 = add_tb(s, I(0.5), I(4.55), I(6), I(0.45))
    first_para(f3, f"{year}  ·  {dept}", 13, False, C_MIDGRAY)

    # Reference sources (bottom)
    rect(s, I(0.07), I(6.85), I(13.26), I(0.02), fill_c=RGBColor(0x2A,0x2A,0x2A))
    f4 = add_tb(s, I(0.5), I(6.9), I(12.5), I(0.38))
    first_para(f4,
        "참고 자료: TalkFile 24/7 무상 기술지원 웹사이트 구축 제안  |  "
        "TalkFile 원격 기술 지원 — 24/7 전문가 대기 체계 제안  |  "
        "최종_선도기업 AM 전략 벤치마킹 보고서",
        8.5, False, RGBColor(0x5A,0x5A,0x5A))
    return s


def slide_toc(prs, chapters):
    s = blank(prs)
    set_bg(s, C_WHITE)

    # Black header bar
    rect(s, I(0), I(0), W_FULL, I(1.05), fill_c=C_BLACK)
    f = add_tb(s, I(0.5), I(0.28), I(4), I(0.5))
    first_para(f, "목  차", 20, True, C_WHITE)

    for i, (num, title, desc) in enumerate(chapters):
        col = i % 2
        row = i // 2
        l = I(0.5  + col * 6.3)
        t = I(1.3  + row * 1.75)

        # Number
        fn = add_tb(s, l, t, I(0.65), I(0.45))
        first_para(fn, num, 12, True, C_CORAL)

        # Title
        ft = add_tb(s, l + I(0.6), t - I(0.02), I(5.4), I(0.48))
        first_para(ft, title, 17, True, C_BLACK)

        # Description
        fd = add_tb(s, l + I(0.6), t + I(0.46), I(5.4), I(0.38))
        first_para(fd, desc, 11, False, C_STEEL)

        # Hairline
        rect(s, l, t + I(0.95), I(5.7), I(0.015), fill_c=C_HAIRLINE)

    # Bottom ref bar
    rect(s, I(0), I(7.1), W_FULL, I(0.4), fill_c=C_SURFACE)
    fb = add_tb(s, I(0.5), I(7.16), I(12.5), I(0.28))
    first_para(fb,
        "참고: TalkFile 24/7 무상 기술지원 웹사이트 구축 제안  |  "
        "TalkFile 원격 기술 지원 — 24/7 전문가 대기 체계 제안  |  "
        "최종_선도기업 AM 전략 벤치마킹 보고서",
        8.5, False, C_STEEL)
    return s


def slide_chapter_div(prs, num, title, subtitle=""):
    s = blank(prs)
    set_bg(s, C_BLACK)
    rect(s, I(0), I(0), I(0.12), H_FULL, fill_c=C_CORAL)

    # Large chapter number
    fn = add_tb(s, I(0.6), I(1.9), I(4), I(1.3))
    first_para(fn, num, 72, True, C_CORAL)

    # Chapter title
    ft = add_tb(s, I(0.6), I(3.25), I(11.5), I(1.0))
    first_para(ft, title, 34, True, C_WHITE, line_spacing=46)

    if subtitle:
        rect(s, I(0.6), I(4.4), I(2.0), I(0.03), fill_c=C_CORAL)
        fs = add_tb(s, I(0.6), I(4.55), I(11), I(0.45))
        first_para(fs, subtitle, 13, False, C_MIDGRAY)
    return s


def slide_content(prs, ch_num, title, items, footer_note=""):
    """
    items: list of tuples
      ('h',  text)        — section heading
      ('b',  text)        — bullet point
      ('sb', text)        — sub-bullet
      ('hl', text)        — coral highlight box (key takeaway)
      ('gap', float)      — vertical gap in inches
    """
    s = blank(prs)
    set_bg(s, C_WHITE)

    # Left coral accent bar (constant)
    rect(s, ACCENT_L, I(0), ACCENT_W, H_FULL, fill_c=C_CORAL)

    # Chapter tag (consistent top-left)
    fn = add_tb(s, TAG_L, TAG_T, I(1.2), I(0.28))
    first_para(fn, ch_num, 9, True, C_CORAL)

    # Slide title (consistent position)
    ft = add_tb(s, TITLE_L, TITLE_T, I(12.5), I(0.72))
    first_para(ft, title, 22, True, C_BLACK, line_spacing=30)

    # Hairline under title (consistent)
    rect(s, ACCENT_W, DIV_T, W_FULL - ACCENT_W, DIV_H, fill_c=C_HAIRLINE)

    # Content rendering
    y = I(1.48)

    def render(item_list):
        nonlocal y
        for item in item_list:
            kind = item[0]

            if kind == 'h':
                y += Pt(3)
                fh = add_tb(s, CONTENT_L, y, CONTENT_W, I(0.46))
                first_para(fh, item[1], 13, True, C_BLACK)
                y += I(0.46)

            elif kind == 'b':
                rect(s, CONTENT_L + I(0.07), y + I(0.12), I(0.055), I(0.055), fill_c=C_CORAL)
                fb = add_tb(s, CONTENT_L + I(0.21), y, CONTENT_W - I(0.22), I(0.43))
                first_para(fb, item[1], 12.5, False, C_CHARCOAL, line_spacing=17.5)
                y += I(0.44)

            elif kind == 'sb':
                fd = add_tb(s, CONTENT_L + I(0.38), y, I(0.18), I(0.38))
                first_para(fd, "–", 11, False, C_STEEL)
                fsb = add_tb(s, CONTENT_L + I(0.58), y, CONTENT_W - I(0.6), I(0.38))
                first_para(fsb, item[1], 11, False, C_STEEL, line_spacing=16)
                y += I(0.39)

            elif kind == 'hl':
                y += I(0.08)
                bh = I(0.68)
                rect(s, CONTENT_L, y, CONTENT_W, bh, fill_c=C_SURFACE, line_c=C_HAIRLINE, line_pt=0.5)
                fhl = add_tb(s, CONTENT_L + I(0.18), y + I(0.12), CONTENT_W - I(0.3), I(0.46))
                first_para(fhl, item[1], 12.5, True, C_CORAL, line_spacing=17)
                y += bh + I(0.05)

            elif kind == 'gap':
                y += I(float(item[1]))

    render(items)

    # Footer note
    if footer_note:
        rect(s, ACCENT_W, I(6.88), W_FULL - ACCENT_W, I(0.022), fill_c=C_HAIRLINE)
        ff = add_tb(s, CONTENT_L, I(6.93), I(12.5), I(0.38))
        first_para(ff, "※ " + footer_note, 9, False, C_STEEL)

    return s


def slide_final(prs, subtitle):
    s = blank(prs)
    set_bg(s, C_BLACK)
    rect(s, I(0), I(0), I(0.07), H_FULL, fill_c=C_CORAL)
    rect(s, I(0), I(0), W_FULL, I(0.05), fill_c=C_CORAL)

    f = add_tb(s, I(0.55), I(2.7), W_FULL - I(1.1), I(1.6))
    first_para(f, "감사합니다", 58, True, C_WHITE, PP_ALIGN.CENTER)

    rect(s, I(5.0), I(4.6), I(3.33), I(0.04), fill_c=C_CORAL)

    f2 = add_tb(s, I(0.55), I(4.75), W_FULL - I(1.1), I(0.5))
    first_para(f2, subtitle, 11, False, C_MIDGRAY, PP_ALIGN.CENTER)
    return s


# ════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ════════════════════════════════════════════════════════════════════
prs = make_prs()

# ── Slide 1: Cover ──────────────────────────────────────────────────
slide_cover(
    prs,
    title_top="회사 중장기('30년까지) 사업계획 달성 및 경쟁력 강화를 위한",
    title_main="웹사이트 구축",
    dept="전략기획팀",
    year="2026"
)

# ── Slide 2: Table of Contents ──────────────────────────────────────
slide_toc(prs, [
    ("01", "오아시스 망한 이유",      "디지털 플랫폼 실패 사례 분석"),
    ("02", "그럼에도 웹사이트?",       "B2B 기술 서비스업의 디지털 채널 필요성"),
    ("03", "선도기업 벤치마킹",        "글로벌 AM 전략 선도기업 사례"),
    ("04", "웹사이트에 넣을 것들",     "CBM 포털 · 24/7 지원 · PM 인계 시스템"),
    ("05", "유상 전환 전략",           "구독형 vs 건별 과금 모델 설계"),
    ("06", "기대효과",                "영업·서비스·수익 구조 개선 성과"),
])

# ── Chapter 01 ──────────────────────────────────────────────────────
slide_chapter_div(prs, "01", "오아시스 망한 이유", "디지털 플랫폼 실패 사례 분석")

slide_content(prs, "01", "오아시스마켓 개요 및 핵심 실패 원인", [
    ('h',  '서비스 개요'),
    ('b',  '새벽배송 기반 신선식품 온라인 플랫폼 (2011년 설립 → 2023년 사실상 경쟁력 상실)'),
    ('b',  '차별화 포인트: 유기농·친환경 상품 중심, 합리적 가격 / 누적 회원 400만 이상 확보'),
    ('b',  '연 거래액 3,000억 규모까지 성장하였으나 흑자 전환 실패 → 경쟁력 급격 약화'),
    ('gap','0.08'),
    ('h',  '핵심 실패 원인'),
    ('b',  '물류 인프라 투자 부담: 콜드체인 유지 비용 급증 → 수익성 구조적 악화'),
    ('sb', '새벽배송 전용 물류센터 구축·운영비 급증, BEP(손익분기점) 도달 불가'),
    ('b',  '경쟁 심화: 쿠팡·컬리·SSG 등 대형 자본 플레이어 진입 → 가격·마케팅 전쟁 감당 불가'),
    ('b',  '수익 모델 취약: 무료배송·과도한 할인쿠폰 → 단위당 수익성 결함 구조화'),
    ('b',  '고객 락인 실패: 반복 구매를 유도할 구독·멤버십 체계 미흡 → 이탈율 지속 증가'),
    ('hl', '→ 디지털 채널 구축 자체가 문제가 아닌, 수익 구조 부재와 차별화 전략 실패가 핵심 원인'),
], footer_note="오아시스마켓은 B2C 신선식품 물류 플랫폼 사례. 당사 B2B 기술 서비스와 사업 구조 상이함")

# ── Chapter 02 ──────────────────────────────────────────────────────
slide_chapter_div(prs, "02", "그럼에도 웹사이트?", "B2B 기술 서비스업에서의 전략적 가치")

slide_content(prs, "02", "우리가 웹사이트를 구축해야 하는 이유", [
    ('h',  'B2B vs B2C: 구조적 차이'),
    ('b',  '오아시스는 불특정 다수 소비자 대상 B2C 물류 플랫폼 → 가격·속도 경쟁에 직면'),
    ('b',  '당사는 선주·선박관리사 대상 B2B 기술 서비스 → 신뢰·전문성·장기 관계가 핵심'),
    ('sb', '고객 수 少, 계약 단가 高, 장기 LTSA 계약 중심 → 플랫폼 경쟁 구조와 무관'),
    ('gap','0.08'),
    ('h',  '웹사이트가 필요한 이유'),
    ('b',  '글로벌 선주·선박관리사는 인터넷으로 AM 파트너를 1차 스크리닝 → 디지털 접점 必'),
    ('b',  '경쟁사 대비 디지털 채널 부재 → 입찰·영업 기회 손실 지속 발생'),
    ('b',  '24/7 기술지원·CBM 데이터 공유를 위한 고객 전용 포털 구축 필요'),
    ('b',  'LTSA 계약 선주 대상 전용 서비스 제공 → 고객 락인 및 재계약율 향상 기여'),
    ('sb', '포털 사용 이력이 쌓일수록 이탈 전환 비용 증가 → 장기 관계 구조적 안정화'),
    ('hl', '→ 웹사이트는 비용이 아닌 영업·서비스·고객 관리를 위한 핵심 인프라 투자'),
])

# ── Chapter 03 ──────────────────────────────────────────────────────
slide_chapter_div(prs, "03", "선도기업 웹사이트 벤치마킹", "글로벌 AM 전략 선도기업 사례 분석")

slide_content(prs, "03", "선도기업 벤치마킹 결과 요약", [
    ('h',  '분석 대상 기업'),
    ('b',  'Wartsila (핀란드): 선박 엔진·추진 시스템 글로벌 1위 / 고객 포털 + CBM 대시보드 운영'),
    ('b',  'MAN Energy Solutions (독일): 저속 엔진 시장 1위 / 24/7 원격 지원 + 디지털 서비스 포털'),
    ('b',  'Kongsberg (노르웨이): 자율운항·디지털 AM 솔루션 선도 / 데이터 기반 예측정비 플랫폼'),
    ('gap','0.08'),
    ('h',  '공통 핵심 특징 및 시사점'),
    ('b',  '고객 전용 포털 운영: 기기별 매뉴얼·부품 조회·서비스 이력 → 고객 접촉 빈도·만족도 향상'),
    ('b',  'CBM 대시보드 탑재: 실시간 상태 모니터링 시각화 → 이상 감지→알림→서비스 연결 자동화'),
    ('b',  '24/7 기술 지원 채널: 챗봇 + 전문가 연결 하이브리드 체계 → 즉각 대응으로 신뢰 확보'),
    ('b',  '케이스 스터디·레퍼런스 콘텐츠 → 신규 고객 유입 및 서비스 전문성 대외 증명'),
    ('hl', '→ 단순 회사 소개를 넘어 고객 서비스 플랫폼으로 진화한 것이 선도기업의 핵심 전략'),
], footer_note="출처: 최종_선도기업 AM 전략 벤치마킹 보고서")

# ── Chapter 04 ──────────────────────────────────────────────────────
slide_chapter_div(prs, "04", "웹사이트에 넣을 것들", "핵심 기능 및 콘텐츠 구성 계획")

slide_content(prs, "04-1", "CBM (Condition Based Monitoring) 웹사이트 탑재", [
    ('h',  'CBM 포털 핵심 기능'),
    ('b',  '실시간 장비 상태 모니터링 대시보드 (선주별 전용 로그인 접근)'),
    ('sb', '진동·온도·압력·오일 분석 데이터 시각화 / 이상 징후 자동 감지 및 알림 (이메일·SMS)'),
    ('b',  '정비 이력·서비스 보고서 디지털 보관 및 조회 (장비·선박 단위 이력 누적)'),
    ('b',  '예측 정비(Predictive Maintenance) 권고 리포트 자동 생성 및 제공'),
    ('gap','0.08'),
    ('h',  'CBM 탑재 기대효과'),
    ('b',  '선주의 장비 운영 가시성 향상 → 당사 서비스 의존도 및 신뢰도 증가'),
    ('b',  'LTSA 계약 갱신율 향상: 데이터로 서비스 가치를 수치로 입증 가능'),
    ('b',  '예방정비 전환으로 긴급 수리 건수·비용 감소 → 고객 만족도 향상'),
    ('sb', '현장 출동 빈도 감소 → 운영 효율화 및 서비스 원가 절감 효과 동시 확보'),
    ('hl', '→ CBM 데이터 포털은 LTSA 계약 차별화 및 신규 선주 유입의 핵심 도구'),
])

slide_content(prs, "04-2", "24/7 무상 기술 지원 서비스", [
    ('h',  '서비스 운영 체계'),
    ('b',  '24시간 365일 기술 지원 전문가 상시 대기 (지역별 로테이션 + 원격 화상 지원 병행)'),
    ('b',  '웹 문의 접수 → 자동 티켓 발행 → 담당자 배정 → 처리 현황 실시간 추적'),
    ('b',  'FAQ + 기술 매뉴얼 자가 검색 기능 → 단순 문의 자동 해결 → 전문가 집중도 향상'),
    ('b',  '원격 진단 도구 연동: 선내 장비 실시간 데이터 연결 → 현장 출동 前 원인 파악'),
    ('gap','0.08'),
    ('h',  '무상 제공 범위 및 유상 전환 기준'),
    ('b',  '무상: 기본 기술 문의·원격 진단·매뉴얼 안내 — LTSA 계약사 전용 혜택'),
    ('b',  '유상: 현장 출동·부품 교체·대규모 개조 공사 — 별도 협의 과금'),
    ('sb', '경계 명확화로 고객 기대 관리 및 안정적 수익 구조를 동시에 확보'),
    ('hl', '→ 무상 기술 지원은 LTSA 신규 계약 유인책이자 고객 충성도 강화의 핵심 수단'),
], footer_note="출처: TalkFile 24/7 무상 기술지원 웹사이트 구축 제안  |  TalkFile 원격 기술 지원 — 24/7 전문가 대기 체계 제안")

slide_content(prs, "04-3", "기술 지원 후 PM 인계 자료 디지털화", [
    ('h',  '현황 및 문제점'),
    ('b',  '현장 기술 지원 완료 후 PM 인계 자료가 이메일·종이 문서로 산재 → 체계 없음'),
    ('b',  '담당자 교체 시 서비스 이력 단절 → 재발 문제 대응 지연 → 고객 불만 증가'),
    ('b',  '선주 자료 요청 시 조회에 과도한 시간 소요 → 대외 전문성 신뢰도 저하'),
    ('gap','0.08'),
    ('h',  '웹사이트 기반 디지털 인계 체계 구축'),
    ('b',  '지원 완료 즉시 포털에 보고서 업로드 (작업 내용·교체 부품·차기 점검일 등 표준화)'),
    ('b',  '선주 및 내부 PM팀 동시 열람 가능 → 정보 투명성 확보 → 신뢰 관계 강화'),
    ('b',  '장비·선박 단위 이력 누적 → 예측 정비 데이터베이스로 전략적 활용'),
    ('b',  'LTSA 연계 활용: 이력 데이터를 신규 선주 영업 레퍼런스로 활용 → 계약 확장 기여'),
    ('hl', '→ LTSA 계약 확대 시 디지털 이력 관리가 추가 선주 유입의 핵심 증거 자료로 기능'),
])

# ── Chapter 05 ──────────────────────────────────────────────────────
slide_chapter_div(prs, "05", "유상전환시점 전략", "무상 → 유상 과금 전환 모델 설계")

slide_content(prs, "05", "유상 전환 시점 및 과금 모델 설계", [
    ('h',  '과금 모델 옵션 비교'),
    ('b',  '[구독형] 월정액 기반 기술 지원 서비스 패키지 — 안정적 반복 매출 확보'),
    ('sb', 'Basic(원격 지원만) / Standard(+현장 연 1회) / Premium(무제한 현장) 3-tier 구조'),
    ('sb', '고객 예산 계획 용이, LTSA 재계약 연동 시 락인 효과 극대화'),
    ('b',  '[건별형] 서비스 발생 건당 단위 과금 — 소규모 선주·단기 계약사 접근성 높음'),
    ('sb', '수익 예측 불안정, 장기 고객 관계 연속성 약함 → 보완 모델로 병행 유지'),
    ('gap','0.08'),
    ('h',  '권장 전환 로드맵'),
    ('b',  "1단계 ('26~'27): 무상 서비스로 LTSA 기존 고객 포털 정착 유도 → 사용 데이터 축적"),
    ('b',  "2단계 ('28): 신규 계약 고객 대상 구독형 Basic 유료화 시작 → 시장 반응 모니터링"),
    ('b',  "3단계 ('29~'30): 전 고객 구독 패키지 전환 완료 / 건별형 병행 유지"),
    ('hl', "→ 구독형 전환 목표: '30년 기준 서비스 반복 매출 비중 30% 이상 달성"),
])

# ── Chapter 06 ──────────────────────────────────────────────────────
slide_chapter_div(prs, "06", "기대효과", "웹사이트 구축을 통한 사업 성과 전망")

slide_content(prs, "06", "웹사이트 구축 기대효과 총괄", [
    ('h',  '영업·마케팅 효과'),
    ('b',  '글로벌 선주·선박관리사 디지털 접점 확보 → 연간 신규 문의 30% 증가 목표'),
    ('b',  '레퍼런스 콘텐츠 + CBM 데이터 시각화로 영업 설득력 강화 → 수주 전환율 향상'),
    ('b',  "LTSA 계약 갱신율 향상 목표: '30년까지 90% 이상 유지"),
    ('gap','0.06'),
    ('h',  '서비스·운영 효과'),
    ('b',  '24/7 원격 지원 체계로 현장 출동 비용 절감 (목표: 연간 출동 건수 20% 감소)'),
    ('b',  '디지털 이력 관리로 PM 인계 소요 시간 50% 단축 → 서비스 효율화'),
    ('b',  'CBM 기반 예측정비 → 긴급 수리 건수 감소 → 고객 만족도·NPS 향상'),
    ('gap','0.06'),
    ('h',  '중장기 수익 구조 개선'),
    ('b',  "구독 서비스 유료화('28년~) → '30년 서비스 반복 매출 비중 30% 달성"),
    ('b',  '웹 포털 → 디지털 AM 플랫폼으로 진화 → 장기 경쟁 우위 구조적 확보'),
    ('hl', "→ 웹사이트 구축은 '30년 중장기 목표 달성을 위한 핵심 디지털 인프라 투자"),
])

# ── Final slide ──────────────────────────────────────────────────────
slide_final(prs, "회사 중장기('30년까지) 사업계획 달성 및 경쟁력 강화를 위한 웹사이트 구축")

# ── Save ─────────────────────────────────────────────────────────────
out_path = '/Users/sinchaeyeon/Desktop/채린/웹사이트구축_PPT.pptx'
prs.save(out_path)
print(f"Saved → {out_path}")
print(f"Total slides: {len(prs.slides)}")
