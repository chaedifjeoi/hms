#!/usr/bin/env python3
"""LTSA 기반 유상 서비스 연계 방향 — PPT 발표용 생성기.

선행 24/7 무상 기술지원 PPT와 동일한 디자인 언어 사용.
약 18장으로 10분 발표 분량 구성.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUTPUT = Path('/Users/sinchaeyeon/Desktop/채린/LTSA_유상_서비스_연계_방향_제안_발표용.pptx')

# Brand palette (same as 선행 PPT)
NAVY        = RGBColor(0x00, 0x2B, 0x5C)
NAVY_DARK   = RGBColor(0x00, 0x1A, 0x3D)
ORANGE      = RGBColor(0xFF, 0x7A, 0x00)
LIGHT_BLUE  = RGBColor(0xE8, 0xEE, 0xF7)
GRAY_TEXT   = RGBColor(0x33, 0x33, 0x33)
GRAY_SUB    = RGBColor(0x70, 0x80, 0x95)
LIGHT_GRAY  = RGBColor(0xCB, 0xD2, 0xD9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0x96, 0xA0)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
RED         = RGBColor(0xC0, 0x39, 0x2B)

KFONT = "Apple SD Gothic Neo"

TOTAL_PAGES = 18  # set after we know how many slides we'll have


# ── shape helpers --------------------------------------------------------

def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=GRAY_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=KFONT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_header(slide, page_num, kicker, title):
    add_rect(slide, Inches(0), Inches(0), Inches(0.3), Inches(1.5), NAVY)
    add_text(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.3),
             kicker, size=11, bold=True, color=ORANGE)
    add_text(slide, Inches(0.5), Inches(0.6), Inches(12), Inches(0.7),
             title, size=26, bold=True, color=NAVY)
    add_rect(slide, Inches(0), Inches(7.4), Inches(13.33), Inches(0.1), ORANGE)
    add_text(slide, Inches(0.5), Inches(7.0), Inches(8), Inches(0.3),
             "HD현대마린솔루션테크 · LTSA 유상 서비스 연계 방향 제안",
             size=9, color=LIGHT_GRAY)
    add_text(slide, Inches(11.5), Inches(7.0), Inches(1.6), Inches(0.3),
             f"{page_num} / {TOTAL_PAGES}",
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def add_blank(prs):
    layout = prs.slide_layouts[6]
    for l in prs.slide_layouts:
        if l.name.lower() == "blank":
            layout = l
            break
    return prs.slides.add_slide(layout)


def make_table(slide, x, y, w, h, headers, rows, *, header_fill=NAVY,
               header_color=WHITE, header_size=11, body_size=10,
               row_alt=LIGHT_BLUE, col_widths=None, first_col_bold=True):
    """Manual table built from rectangles + textboxes for full styling control."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +header
    if col_widths is None:
        col_widths = [w / n_cols] * n_cols
    else:
        col_widths = [Emu(int(w * c / sum(col_widths))) for c in col_widths]
    row_h = Emu(int(h / n_rows))

    # header row
    cur_x = x
    for i, head in enumerate(headers):
        add_rect(slide, cur_x, y, col_widths[i], row_h, header_fill)
        add_text(slide, cur_x + Inches(0.1), y, col_widths[i] - Inches(0.2), row_h,
                 head, size=header_size, bold=True, color=header_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x = cur_x + col_widths[i]

    # body rows
    for ri, row in enumerate(rows):
        cur_x = x
        ry = y + row_h * (ri + 1)
        if ri % 2 == 0:
            add_rect(slide, x, ry, w, row_h, row_alt)
        for ci, val in enumerate(row):
            bold = first_col_bold and ci == 0
            color = NAVY if bold else GRAY_TEXT
            # detect highlight markers
            txt = str(val)
            if txt.startswith("**") and txt.endswith("**"):
                txt = txt[2:-2]
                bold = True
                color = ORANGE
            add_text(slide, cur_x + Inches(0.1), ry, col_widths[ci] - Inches(0.2),
                     row_h, txt, size=body_size, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
            cur_x = cur_x + col_widths[ci]


# ── slides ---------------------------------------------------------------

def slide_title(prs):
    s = add_blank(prs)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), NAVY)
    add_rect(s, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7), NAVY_DARK)
    add_rect(s, Inches(11.5), Inches(0), Inches(1.83), Inches(0.5), ORANGE)

    add_text(s, Inches(0.7), Inches(1.0), Inches(8), Inches(0.4),
             "EXECUTIVE PROPOSAL · 2026.04", size=13, bold=True, color=ORANGE)
    add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(1.2),
             "LTSA 기반 유상 서비스 연계 방향",
             size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(3.4), Inches(12), Inches(0.8),
             "2030년 안정 매출 구조와 업무 프로세스 개선 통합 제안",
             size=22, color=WHITE)
    add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.5),
             "무상 기술지원 → AI 케이스 인계 → HI-CMS 알림 → LTSA 자연 전환",
             size=14, color=LIGHT_GRAY)

    add_rect(s, Inches(0.7), Inches(5.4), Inches(0.5), Inches(0.06), ORANGE)
    add_text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.4),
             "회사 주제 통합 답변 · 중장기 매출 증대 + 업무 프로세스 개선",
             size=12, color=LIGHT_GRAY, italic=True)

    add_text(s, Inches(0.7), Inches(7.0), Inches(8), Inches(0.3),
             "HD현대마린솔루션테크", size=12, bold=True, color=WHITE)
    add_text(s, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
             "2026. 04", size=11, color=ORANGE, align=PP_ALIGN.RIGHT)


def slide_executive(prs):
    s = add_blank(prs)
    add_header(s, 2, "EXECUTIVE SUMMARY", "한 장으로 보는 제안")

    # 핵심 명제 box
    add_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), LIGHT_BLUE)
    add_rect(s, Inches(0.5), Inches(1.5), Inches(0.12), Inches(1.2), ORANGE)
    add_text(s, Inches(0.85), Inches(1.6), Inches(12), Inches(0.3),
             "핵심 명제", size=11, bold=True, color=ORANGE)
    add_text(s, Inches(0.85), Inches(1.85), Inches(12), Inches(0.4),
             '무상 보증기간에 축적되는 기술지원 이력을 자산화해, LTSA(장기 서비스 계약) 안정 매출로 자연 전환한다.',
             size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.85), Inches(2.25), Inches(12), Inches(0.4),
             '매출 증대(중장기 계획)는 결과, 업무 프로세스 개선(AI 인계)은 그 과정에서 부수 달성.',
             size=11, color=GRAY_TEXT)

    # 4축 카드
    cards = [
        ("입구",   "무상 기술지원 포털",   "메일·전화에 산재", "전 선주 단일 창구",     ORANGE),
        ("통로 1", "AI 케이스 인계",      "PM 수기 정리 3~8일",  "자동 요약 15~30분",      TEAL),
        ("통로 2", "HI-CMS 통합 알림",    "사내 분석 한정",      "선주 직접 알림 클릭",     NAVY),
        ("출구",   "LTSA 자연 전환",      "단발 유상 ±25%",      "구독+사건 ±5%",          ORANGE),
    ]
    card_w = Inches(3.0)
    for i, (kicker, title, before, after, color) in enumerate(cards):
        x = Inches(0.5) + i * (card_w + Inches(0.1))
        y = Inches(2.95)
        add_rect(s, x, y, card_w, Inches(2.3), WHITE, line=LIGHT_GRAY)
        add_rect(s, x, y, card_w, Inches(0.4), color)
        add_text(s, x + Inches(0.15), y + Inches(0.05),
                 card_w - Inches(0.3), Inches(0.3), kicker,
                 size=11, bold=True, color=WHITE)
        add_text(s, x + Inches(0.2), y + Inches(0.55),
                 card_w - Inches(0.4), Inches(0.5), title,
                 size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), y + Inches(1.1),
                 card_w - Inches(0.4), Inches(0.3), "현재", size=9, color=GRAY_SUB)
        add_text(s, x + Inches(0.2), y + Inches(1.32),
                 card_w - Inches(0.4), Inches(0.4), before, size=10, color=GRAY_TEXT)
        add_text(s, x + Inches(0.2), y + Inches(1.7),
                 card_w - Inches(0.4), Inches(0.3), "2030 목표", size=9, color=ORANGE, bold=True)
        add_text(s, x + Inches(0.2), y + Inches(1.92),
                 card_w - Inches(0.4), Inches(0.4), after, size=11, bold=True, color=NAVY)

    # bottom callout
    add_rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.7), NAVY)
    add_text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.3),
             "회사 주제 통합 답변", size=10, bold=True, color=ORANGE)
    add_text(s, Inches(0.7), Inches(5.78), Inches(12), Inches(0.4),
             "2030 매출 증대(중장기 계획)  +  업무 프로세스 개선  →  단일 메커니즘으로 동시 달성",
             size=14, bold=True, color=WHITE)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
             "※ 본 제안은 「24/7 무상 기술지원 웹사이트 구축 제안」(선행, 2026.04)의 후속이며, 구축된 무상 포털 위에 단계적으로 적용 가능.",
             size=10, italic=True, color=GRAY_SUB)


def slide_toc(prs):
    s = add_blank(prs)
    add_header(s, 3, "TABLE OF CONTENTS", "본 발표가 답할 5개 질문")

    qs = [
        ("Q1", "왜 LTSA인가?",                    "OEM 3사가 30년 걸려 도달한 매출 안정 구조의 본질",     "본론 1"),
        ("Q2", "무상 → 유상 단절은 어떻게 해소하나?", "무상 포털 자체가 LTSA 영업 도구가 되는 구조",            "본론 2"),
        ("Q3", "AI 케이스 인계는 어떻게 작동하나?",   "PM 수기 3~8일 → AI 자동 15~30분 (99% 감축)",           "본론 3"),
        ("Q4", "HI-CMS는 어떻게 선주를 끌어오나?",    "이상 감지→웹 알림→케이스 자동 생성 (선주 클릭 1회)",   "본론 4"),
        ("Q5", "2030년까지 어떻게 도달하나?",         "5-Phase 게이트 / LTSA 비중 5% → 50%↑",                "본론 5"),
    ]
    y0 = Inches(1.6)
    for i, (qk, qt, ans, sec) in enumerate(qs):
        y = y0 + i * Inches(1.05)
        add_rect(s, Inches(0.5), y, Inches(0.9), Inches(0.9), NAVY)
        add_text(s, Inches(0.5), y, Inches(0.9), Inches(0.9),
                 qk, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        add_rect(s, Inches(1.4), y, Inches(11.4), Inches(0.9), LIGHT_BLUE)
        add_text(s, Inches(1.6), y + Inches(0.1), Inches(8.5), Inches(0.4),
                 qt, size=14, bold=True, color=NAVY)
        add_text(s, Inches(1.6), y + Inches(0.5), Inches(8.5), Inches(0.4),
                 ans, size=11, color=GRAY_TEXT)
        add_text(s, Inches(11.5), y + Inches(0.3), Inches(1.2), Inches(0.4),
                 sec, size=11, bold=True, color=ORANGE,
                 align=PP_ALIGN.RIGHT)


def slide_intro(prs):
    s = add_blank(prs)
    add_header(s, 4, "서론", "회사 주제 두 가지를 하나의 메커니즘으로")

    # two-column comparison: 회사 주제 둘 → 통합
    add_rect(s, Inches(0.5), Inches(1.6), Inches(5.7), Inches(2.0), LIGHT_BLUE)
    add_rect(s, Inches(0.5), Inches(1.6), Inches(0.12), Inches(2.0), GRAY_SUB)
    add_text(s, Inches(0.85), Inches(1.7), Inches(5.5), Inches(0.3),
             "주제 A · 중장기 매출 증대", size=11, bold=True, color=GRAY_SUB)
    add_text(s, Inches(0.85), Inches(2.0), Inches(5.5), Inches(0.5),
             "2030년까지 안정적 매출 구조 확보", size=18, bold=True, color=NAVY)
    add_text(s, Inches(0.85), Inches(2.6), Inches(5.5), Inches(1.0),
             "▸ LTSA 매출 비중 5~10% → 50%↑\n▸ 매출 예측 정확도 ±25% → ±5%\n▸ 단발성 견적 → 안정 구독 구조",
             size=11, color=GRAY_TEXT)

    add_rect(s, Inches(7.13), Inches(1.6), Inches(5.7), Inches(2.0), LIGHT_BLUE)
    add_rect(s, Inches(7.13), Inches(1.6), Inches(0.12), Inches(2.0), GRAY_SUB)
    add_text(s, Inches(7.5), Inches(1.7), Inches(5.5), Inches(0.3),
             "주제 B · 업무 프로세스 개선", size=11, bold=True, color=GRAY_SUB)
    add_text(s, Inches(7.5), Inches(2.0), Inches(5.5), Inches(0.5),
             "PM·기술지원·HI-CMS 운영 효율화", size=18, bold=True, color=NAVY)
    add_text(s, Inches(7.5), Inches(2.6), Inches(5.5), Inches(1.0),
             "▸ PM 1척당 인계 3~8일 → 15~30분\n▸ HI-CMS 사내 한정 → 선주 직접 도달\n▸ R&R·KPI 측정 가능 구조",
             size=11, color=GRAY_TEXT)

    # arrow merging downward
    arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                              Inches(6.16), Inches(3.7),
                              Inches(1.0), Inches(0.5))
    arr.fill.solid(); arr.fill.fore_color.rgb = ORANGE
    arr.line.fill.background()

    # combined card
    add_rect(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5), NAVY)
    add_rect(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5), ORANGE)
    add_text(s, Inches(0.85), Inches(4.4), Inches(12), Inches(0.3),
             "본 제안의 통합 답변", size=11, bold=True, color=WHITE)
    add_text(s, Inches(0.85), Inches(4.95), Inches(12), Inches(0.6),
             '"무상 기술지원 → AI 케이스 인계 → HI-CMS 알림 → LTSA 자연 전환"',
             size=20, bold=True, color=WHITE)
    add_text(s, Inches(0.85), Inches(5.55), Inches(12), Inches(0.4),
             "두 주제는 별도 답이 아니라 하나의 데이터 흐름에서 동시 달성된다.",
             size=12, color=LIGHT_GRAY)
    add_text(s, Inches(0.85), Inches(6.0), Inches(12), Inches(0.7),
             "중장기 매출 증대(=결과)  ·  업무 프로세스 개선(=과정)  ·  신규 영업·신규 인력 일괄 투입 불필요",
             size=12, bold=True, color=ORANGE)


def slide_status(prs):
    s = add_blank(prs)
    add_header(s, 5, "본론 1 · Q1 — 현황과 단절",
               "무상-유상 사이의 데이터 단절 + PM 케이스 재구성 부담")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "무상 보증기간 동안 쌓인 기술지원 이력이 유상 서비스 조직(PM)에 도달하지 않음 — "
             "구조적 단절이 매출 손실의 원인",
             size=12, color=GRAY_SUB)

    # PM 5단계 process - flow boxes
    steps = [
        ("1", "담당 엔지니어 식별",       "0.5~1일"),
        ("2", "메일 이력 수집",          "1~3일"),
        ("3", "케이스 정리·재구성",       "0.5~2일"),
        ("4", "선박 데이터 별도 조회",     "0.5~1일"),
        ("5", "유상 견적·범위 산정",       "0.5~1일"),
    ]
    box_w = Inches(2.3)
    gap = Inches(0.15)
    start_x = Inches(0.5)
    y = Inches(2.0)

    for i, (n, label, time) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, y, box_w, Inches(1.7), LIGHT_BLUE)
        add_rect(s, x, y, box_w, Inches(0.5), NAVY)
        add_text(s, x, y + Inches(0.05), box_w, Inches(0.4),
                 f"단계 {n}", size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), y + Inches(0.65),
                 box_w - Inches(0.2), Inches(0.6), label,
                 size=11, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), y + Inches(1.25),
                 box_w - Inches(0.2), Inches(0.4), time,
                 size=12, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + box_w + Emu(int(gap * 0.05)),
                                      y + Inches(0.7), gap, Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()

    # total time bar
    add_rect(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.7), NAVY)
    add_text(s, Inches(0.7), Inches(4.05), Inches(7), Inches(0.3),
             "PM 1척당 케이스 재구성 총 소요 시간", size=11, bold=True, color=ORANGE)
    add_text(s, Inches(0.7), Inches(4.3), Inches(7), Inches(0.4),
             "수기 메일 정리 기반 → 매번 처음부터 반복",
             size=12, color=LIGHT_GRAY)
    add_text(s, Inches(8.5), Inches(4.05), Inches(4.0), Inches(0.6),
             "3 ~ 8일 / 척", size=26, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # consequence
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
             "결과 — 인계 지연 = LTSA 전환 기회 손실",
             size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5),
             "▸ PM 1인이 동시 처리 가능한 선박 수 제한 (20~30척)\n"
             "▸ LTSA 제안 발송 평균 2주~1개월 지연 → 선주는 그 사이 단발 견적으로 종결\n"
             "▸ 영업 기회 손실률 30~50% (지연 기인) — 인력 추가만으론 근본 해결 불가",
             size=12, color=GRAY_TEXT)


def slide_revenue_gap(prs):
    s = add_blank(prs)
    add_header(s, 6, "본론 1 · Q1 — 매출 안정성 갭",
               "OEM 3사 LCA 매출 비중 vs HMS 현재")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "OEM 3사는 매출의 40~60%가 LCA(LTSA) — 안정·예측 가능. HMS는 5~10%로 단발 의존.",
             size=12, color=GRAY_SUB)

    # comparison bars
    bar_left = Inches(2.5)
    bar_w_max = Inches(8.5)
    rows = [
        ("OEM 3사 평균",        "40~60%",  0.50,  NAVY),
        ("HMS 현재 (추정)",     "5~10%",   0.075, GRAY_SUB),
        ("HMS 2030 목표 (본 제안)", "50%↑",  0.50,  ORANGE),
    ]
    bar_y = Inches(2.2)
    for i, (label, val, ratio, color) in enumerate(rows):
        y = bar_y + i * Inches(1.0)
        add_text(s, Inches(0.5), y, Inches(2.0), Inches(0.4),
                 label, size=12, bold=True, color=NAVY)
        add_text(s, Inches(0.5), y + Inches(0.32), Inches(2.0), Inches(0.3),
                 "LCA/LTSA 매출 비중", size=9, color=GRAY_SUB)
        add_rect(s, bar_left, y + Inches(0.05), bar_w_max, Inches(0.6), LIGHT_BLUE)
        add_rect(s, bar_left, y + Inches(0.05),
                 Emu(int(bar_w_max * ratio)), Inches(0.6), color)
        add_text(s, bar_left + Emu(int(bar_w_max * ratio)) + Inches(0.15),
                 y + Inches(0.1), Inches(2.5), Inches(0.5),
                 val, size=22, bold=True, color=color,
                 anchor=MSO_ANCHOR.MIDDLE)

    # gap callout
    add_rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.7), LIGHT_BLUE)
    add_rect(s, Inches(0.5), Inches(5.5), Inches(0.12), Inches(0.7), ORANGE)
    add_text(s, Inches(0.85), Inches(5.55), Inches(12), Inches(0.3),
             "지표 비교 (OEM 3사 평균 → HMS 2030 목표)",
             size=10, bold=True, color=ORANGE)
    add_text(s, Inches(0.85), Inches(5.8), Inches(12), Inches(0.4),
             "매출 예측 정확도  ±5%   ·   장기 고객 이탈률  <5%   ·   선박당 연 LCA 수익  확보",
             size=12, bold=True, color=NAVY)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
             "→ 2030년까지 LTSA 비중을 OEM 3사 수준으로 끌어올리는 것이 본 제안의 핵심 골자.",
             size=12, italic=True, color=GRAY_SUB)


def slide_ltsa_flow(prs):
    s = add_blank(prs)
    add_header(s, 7, "본론 2 · Q2 — LTSA 연계 흐름",
               "무상 입구 → 데이터 자산 → 유상 출구 (4단계)")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "선주의 무상 포털 사용 자체가 LTSA의 영업 도구가 됨 — 별도 영업 캠페인 불필요",
             size=12, color=GRAY_SUB)

    stages = [
        ("Stage 1", "도입",      "인도 ~ 보증 1년차",   "무상 포털 사용 시작",        "데이터 자산 축적"),
        ("Stage 2", "신뢰",      "보증 1~2년차",       "응답 속도·이력 일원화 체감",   "AI 케이스 인계 가동"),
        ("Stage 3", "전환 제안", "보증 종료 6개월 전",  "LTSA 옵션 자동 안내",         "예측 매출 사전 편입"),
        ("Stage 4", "운영",      "보증 종료 후 N년",    "구독+사건 LTSA 사용",         "갱신 유도·매출 안정화"),
    ]
    box_w = Inches(2.85)
    gap = Inches(0.2)
    y = Inches(2.0)
    for i, (st, name, period, owner, hms) in enumerate(stages):
        x = Inches(0.5) + i * (box_w + gap)
        add_rect(s, x, y, box_w, Inches(3.6), LIGHT_BLUE)
        add_rect(s, x, y, box_w, Inches(0.5), NAVY if i % 2 == 0 else NAVY_DARK)
        add_text(s, x, y + Inches(0.05), box_w, Inches(0.4),
                 st, size=11, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), y + Inches(0.6), box_w - Inches(0.2),
                 Inches(0.5), name, size=18, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), y + Inches(1.15), box_w - Inches(0.2),
                 Inches(0.3), period, size=10, color=ORANGE, bold=True,
                 align=PP_ALIGN.CENTER)
        # divider
        add_rect(s, x + Inches(0.3), y + Inches(1.55),
                 box_w - Inches(0.6), Emu(1), LIGHT_GRAY)

        add_text(s, x + Inches(0.15), y + Inches(1.7), box_w - Inches(0.3),
                 Inches(0.3), "선주 입장", size=9, color=GRAY_SUB, bold=True)
        add_text(s, x + Inches(0.15), y + Inches(1.95), box_w - Inches(0.3),
                 Inches(0.7), owner, size=10, color=GRAY_TEXT)

        add_text(s, x + Inches(0.15), y + Inches(2.65), box_w - Inches(0.3),
                 Inches(0.3), "HMS 입장", size=9, color=GRAY_SUB, bold=True)
        add_text(s, x + Inches(0.15), y + Inches(2.9), box_w - Inches(0.3),
                 Inches(0.7), hms, size=10, color=NAVY, bold=True)

        if i < len(stages) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + box_w + Emu(int(gap * 0.1)),
                                      y + Inches(1.5), gap, Inches(0.4))
            arr.fill.solid(); arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()

    add_rect(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.7), NAVY)
    add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.3),
             "전환 트리거", size=10, bold=True, color=ORANGE)
    add_text(s, Inches(0.7), Inches(6.18), Inches(12), Inches(0.4),
             "보증 종료 · 고빈도 기술 문의 · HI-CMS 이상 감지 · 대형 케이스 발생 · 신규 선박 인도 — 5종이 활성화될 때 자동 LTSA 제안",
             size=11, color=WHITE)

    add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.3),
             "선주는 \"가입 결정\"을 하지 않음 — 자기 사용 패턴이 자동으로 LTSA 제안을 트리거함.",
             size=10, italic=True, color=GRAY_SUB)


def slide_ltsa_models(prs):
    s = add_blank(prs)
    add_header(s, 8, "본론 2 · Q2 — LTSA 모델 비교",
               "구독형 / 사건 기반 / 하이브리드 — 선주·HMS 양쪽 효율")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "결론: 하이브리드(기본료+초과 사용량) 모델이 양쪽 효율 모두 최대 — 본 제안의 권고",
             size=12, color=GRAY_SUB)

    # 3 model cards
    models = [
        ("A", "구독형",     "Subscription",    "연 정액 (선박당)",
         "예측 가능 / 사용량 무관 부담",     "매우 안정 (±5%)",   "소형 선주 적합", GRAY_SUB),
        ("B", "사건 기반",  "Pay-per-Service", "발생 시 청구",
         "저사용 시 유리 / 변동 큼",         "불안정 (±25%)",     "대형 자가정비 선주", NAVY),
        ("C", "하이브리드", "Hybrid (권고)",   "기본료 + 초과 사용",
         "예측 가능 + 사용량 반영",          "안정 (±10%)",       "전 선주 기본 권고", ORANGE),
    ]
    card_w = Inches(4.05)
    gap = Inches(0.1)
    y = Inches(2.0)
    for i, (idx, name, en, structure, owner, stable, fit, color) in enumerate(models):
        x = Inches(0.5) + i * (card_w + gap)
        add_rect(s, x, y, card_w, Inches(4.0), WHITE, line=LIGHT_GRAY)
        add_rect(s, x, y, card_w, Inches(0.6), color)
        add_text(s, x + Inches(0.2), y + Inches(0.07), Inches(0.6), Inches(0.45),
                 idx, size=22, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), y + Inches(0.05), card_w - Inches(0.95), Inches(0.3),
                 name, size=15, bold=True, color=WHITE)
        add_text(s, x + Inches(0.85), y + Inches(0.32), card_w - Inches(0.95), Inches(0.3),
                 en, size=10, color=WHITE)

        rows = [
            ("구조",            structure),
            ("선주 부담",        owner),
            ("HMS 매출 안정성",   stable),
            ("적합 대상",         fit),
        ]
        for ri, (k, v) in enumerate(rows):
            ry = y + Inches(0.85) + ri * Inches(0.75)
            add_text(s, x + Inches(0.2), ry, card_w - Inches(0.4), Inches(0.3),
                     k, size=10, color=GRAY_SUB, bold=True)
            add_text(s, x + Inches(0.2), ry + Inches(0.27), card_w - Inches(0.4),
                     Inches(0.4), v, size=11,
                     bold=(idx == "C"),
                     color=ORANGE if idx == "C" else NAVY)

    add_rect(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), LIGHT_BLUE)
    add_rect(s, Inches(0.5), Inches(6.2), Inches(0.12), Inches(0.7), ORANGE)
    add_text(s, Inches(0.85), Inches(6.25), Inches(12), Inches(0.3),
             "권고 한 줄", size=10, bold=True, color=ORANGE)
    add_text(s, Inches(0.85), Inches(6.5), Inches(12), Inches(0.4),
             "기본료로 매출 안정성 확보 + 실제 사용량(원격진단·부품·엔지니어 파견)은 별도 정산 — 양쪽 효율 동시 충족",
             size=12, bold=True, color=NAVY)


def slide_ai_flow(prs):
    s = add_blank(prs)
    add_header(s, 9, "본론 3 · Q3 — AI 케이스 인계",
               "케이스 종료 시점, AI가 PM에게 자동 인계")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "PM이 메일을 뒤지지 않고도 LTSA 제안에 필요한 정보를 즉시 받음",
             size=12, color=GRAY_SUB)

    # flow with 5 steps
    steps = [
        ("1", "케이스 종료",     "기술지원 엔지니어",    "0초"),
        ("2", "AI 요약 생성",    "AI 에이전트",        "5~10초"),
        ("3", "PM 대시보드 푸시", "시스템",            "실시간"),
        ("4", "PM 검토",        "PM",                "5~15분"),
        ("5", "LTSA 제안 발송",  "PM (자동 템플릿)",   "실시간"),
    ]
    box_w = Inches(2.3)
    gap = Inches(0.18)
    y = Inches(2.0)
    for i, (n, label, owner, t) in enumerate(steps):
        x = Inches(0.5) + i * (box_w + gap)
        add_rect(s, x, y, box_w, Inches(1.85), WHITE, line=LIGHT_GRAY)
        add_rect(s, x, y, Inches(0.5), Inches(1.85), ORANGE)
        add_text(s, x, y, Inches(0.5), Inches(1.85),
                 n, size=24, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.6), y + Inches(0.15),
                 box_w - Inches(0.7), Inches(0.5), label,
                 size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.6), y + Inches(0.85),
                 box_w - Inches(0.7), Inches(0.3), owner, size=9, color=GRAY_SUB)
        add_text(s, x + Inches(0.6), y + Inches(1.2),
                 box_w - Inches(0.7), Inches(0.5), t, size=14, bold=True, color=ORANGE)

        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + box_w, y + Inches(0.8), gap, Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
            arr.line.fill.background()

    # AI 요약 7개 항목
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.4),
             "AI 요약이 자동 생성하는 7개 항목",
             size=13, bold=True, color=NAVY)
    items = [
        "① 장비·부위 (계통 분류)",
        "② 증상 요약",
        "③ 진단 결과 (Tier 1~4 종합)",
        "④ 임시 조치 내역",
        "⑤ 잔여 리스크 평가",
        "⑥ LTSA 권고 액션",
        "⑦ 매출 기회 추정",
    ]
    item_w = Inches(1.7)
    iy = Inches(4.6)
    for i, it in enumerate(items):
        x = Inches(0.5) + i * (item_w + Inches(0.05))
        add_rect(s, x, iy, item_w, Inches(0.55), LIGHT_BLUE)
        add_text(s, x + Inches(0.1), iy, item_w - Inches(0.2), Inches(0.55),
                 it, size=10, bold=True, color=NAVY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # bottom highlight
    add_rect(s, Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.95), NAVY)
    add_text(s, Inches(0.7), Inches(5.55), Inches(8), Inches(0.4),
             "총 소요 시간",
             size=11, bold=True, color=ORANGE)
    add_text(s, Inches(0.7), Inches(5.85), Inches(7), Inches(0.5),
             "수기 PM 정리 → 3 ~ 8일",
             size=14, color=LIGHT_GRAY)
    add_text(s, Inches(7.5), Inches(5.55), Inches(5.3), Inches(0.4),
             "AI 인계 적용 후",
             size=11, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(7.5), Inches(5.85), Inches(5.3), Inches(0.5),
             "15 ~ 30분 (99% 감축)",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             "→ PM 1인당 동시 처리 가능 선박 5~7배 증가  ·  LTSA 제안 발송 지연 평균 2주 → 실시간",
             size=11, italic=True, color=GRAY_SUB)


def slide_ai_impact(prs):
    s = add_blank(prs)
    add_header(s, 10, "본론 3 · Q3 — AI 인계 정량 효과",
               "PM 처리 능력과 매출 기회의 동시 점프")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "Before(수기 인계) vs After(AI 인계) — 6개 핵심 지표 비교",
             size=12, color=GRAY_SUB)

    metrics = [
        ("PM 1인당 동시 처리 선박",   "20~30척",         0.15,  "150~200척",       1.00),
        ("케이스 인계 소요",          "3~8일",           0.95,  "15~30분",         0.04),
        ("LTSA 제안 발송 지연",       "평균 2주",        1.00,  "실시간",            0.02),
        ("PM 1인당 연 LTSA 제안",     "40~60건",         0.15,  "300~400건",       1.00),
        ("AI 요약 신뢰도",            "—",              0.0,   "95%+ (검토 5분)",  0.95),
        ("영업 기회 손실률",          "30~50%",          1.00,  "5% 이하",          0.05),
    ]
    chart_top = Inches(1.95)
    row_h = Inches(0.55)
    gap_y = Inches(0.07)
    label_w = Inches(3.2)
    bar_zone_x = Inches(3.9)
    bar_zone_w = Inches(6.2)
    text_zone_x = bar_zone_x + bar_zone_w + Inches(0.1)
    text_zone_w = Inches(2.9)

    # legend
    lg_y = Inches(1.45)
    add_rect(s, Inches(9.0), lg_y, Inches(0.3), Inches(0.18), GRAY_SUB)
    add_text(s, Inches(9.35), lg_y - Inches(0.02), Inches(1.3), Inches(0.25),
             "Before (수기)", size=10, color=GRAY_TEXT)
    add_rect(s, Inches(11.0), lg_y, Inches(0.3), Inches(0.18), ORANGE)
    add_text(s, Inches(11.35), lg_y - Inches(0.02), Inches(1.5), Inches(0.25),
             "After (AI 인계)", size=10, color=GRAY_TEXT)

    for i, (label, b_text, b_pct, a_text, a_pct) in enumerate(metrics):
        y = chart_top + i * (row_h + gap_y)
        if i % 2 == 0:
            add_rect(s, Inches(0.5), y - Inches(0.03),
                     Inches(12.3), row_h + Inches(0.06), LIGHT_BLUE)
        add_text(s, Inches(0.6), y + Inches(0.13), label_w, Inches(0.3),
                 label, size=11, bold=True, color=NAVY)

        before_w = Emu(int(bar_zone_w * 0.45 * b_pct))
        add_rect(s, bar_zone_x, y + Inches(0.04),
                 max(before_w, Emu(1)), Inches(0.18), GRAY_SUB)
        after_w = Emu(int(bar_zone_w * a_pct))
        add_rect(s, bar_zone_x, y + Inches(0.27),
                 max(after_w, Emu(1)), Inches(0.20), ORANGE)
        add_text(s, text_zone_x, y + Inches(0.0), text_zone_w, Inches(0.25),
                 b_text, size=10, color=GRAY_SUB)
        add_text(s, text_zone_x, y + Inches(0.25), text_zone_w, Inches(0.25),
                 a_text, size=10, bold=True, color=ORANGE)

    add_rect(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.7), NAVY)
    add_rect(s, Inches(0.5), Inches(6.15), Inches(0.12), Inches(0.7), ORANGE)
    add_text(s, Inches(0.85), Inches(6.20), Inches(12), Inches(0.3),
             "정량 효과 한 줄 요약",
             size=10, bold=True, color=ORANGE)
    add_text(s, Inches(0.85), Inches(6.45), Inches(12), Inches(0.4),
             "처리 능력 ×5~7   ·   인계 시간 ↓99%   ·   제안 지연 ↓99%   ·   영업 기회 손실 1/10",
             size=14, bold=True, color=WHITE)


def slide_hicms(prs):
    s = add_blank(prs)
    add_header(s, 11, "본론 4 · Q4 — HI-CMS 통합 알림",
               "선주를 \"와야 하는 사람\"으로 전환하는 자석")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "HI-CMS = HI-ODS(운영 데이터) + ISS(통합 스마트 선박) — 풍부한 데이터를 선주에게 가공 알림으로 도달",
             size=12, color=GRAY_SUB)

    # 4-step flow
    steps = [
        ("1", "이상 감지",     "HI-CMS",      "임계치 초과 / 학습 모델",   "(선주는 모름)"),
        ("2", "웹 알림",      "포털 + 모바일", "푸시 + 메일",              "선주 알림 클릭"),
        ("3", "케이스 자동 생성", "시스템",      "데이터·이력 자동 첨부",     "기술지원 즉시 시작"),
        ("4", "LTSA 트리거",   "AI 인계",      "예측정비 LTSA 제안",       "선주 클릭 1회로 진행"),
    ]
    box_w = Inches(2.95)
    gap = Inches(0.15)
    y = Inches(1.85)
    box_h = Inches(2.0)
    for i, (n, name, sys, sub, owner) in enumerate(steps):
        x = Inches(0.5) + i * (box_w + gap)
        add_rect(s, x, y, box_w, box_h, LIGHT_BLUE)
        add_rect(s, x, y, box_w, Inches(0.45), NAVY)
        add_text(s, x, y + Inches(0.04), box_w, Inches(0.37),
                 f"단계 {n}", size=10, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), y + Inches(0.5),
                 box_w - Inches(0.3), Inches(0.45), name,
                 size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(1.0),
                 box_w - Inches(0.3), Inches(0.3), sys,
                 size=10, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(1.28),
                 box_w - Inches(0.3), Inches(0.30), sub,
                 size=9, color=GRAY_TEXT,
                 align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.3), y + Inches(1.62),
                 box_w - Inches(0.6), Emu(1), LIGHT_GRAY)
        add_text(s, x + Inches(0.15), y + Inches(1.66),
                 box_w - Inches(0.3), Inches(0.30), "← " + owner,
                 size=9, italic=True, color=NAVY,
                 align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + box_w, y + Inches(0.85),
                                      gap, Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()

    # 5 trigger examples
    add_text(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.4),
             "알림 트리거 예시 (실제 적용 가능 시나리오)",
             size=13, bold=True, color=NAVY)

    triggers = [
        ("연료 효율 저하",   "주행/연료 5%↓ 지속",     "정기 정비 LTSA"),
        ("엔진 진동 증가",   "진동 +15% 이상",          "예측정비 LTSA"),
        ("윤활유 이상",      "오일 분석 + 압력",         "소모품 LTSA"),
        ("DPS 결함 신호",   "자가진단 코드",            "긴급 LTSA + 부품"),
        ("배기·NOx 임계",  "센서 + 운항 패턴",          "규제 대응 LTSA"),
    ]
    tw = Inches(2.45)
    ty = Inches(4.5)
    for i, (name, signal, ltsa) in enumerate(triggers):
        x = Inches(0.5) + i * (tw + Inches(0.05))
        add_rect(s, x, ty, tw, Inches(1.5), WHITE, line=LIGHT_GRAY)
        add_rect(s, x, ty, tw, Inches(0.45), ORANGE)
        add_text(s, x + Inches(0.1), ty + Inches(0.05),
                 tw - Inches(0.2), Inches(0.35), name,
                 size=11, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), ty + Inches(0.55),
                 tw - Inches(0.3), Inches(0.3), "감지 신호",
                 size=8, color=GRAY_SUB)
        add_text(s, x + Inches(0.15), ty + Inches(0.78),
                 tw - Inches(0.3), Inches(0.3), signal,
                 size=10, color=GRAY_TEXT)
        add_text(s, x + Inches(0.15), ty + Inches(1.08),
                 tw - Inches(0.3), Inches(0.3), "→ " + ltsa,
                 size=10, bold=True, color=NAVY)

    add_text(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.5),
             "핵심: HI-CMS 알림은 \"선주가 안 와도 됐던 영역\"을 \"즉시 확인하게 만드는 영역\"으로 바꾼다.",
             size=11, bold=True, italic=True, color=NAVY)


def slide_roadmap(prs):
    s = add_blank(prs)
    add_header(s, 12, "본론 5 · Q5 — 2030 중장기 로드맵",
               "5-Phase 단계별 게이트와 KPI")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "PoC(2026) → 정착(2027) → 확장(2028) → 통합(2029) → 안정(2030)",
             size=12, color=GRAY_SUB)

    phases = [
        ("Phase 1", "MVP",     "2026 1H",      "무상 웹 + AI PoC + 알림 1종",     "10~20개사 Beta",         NAVY),
        ("Phase 2", "정착",    "2026 2H",       "AI 신뢰도 95% / 알림 5종 / LTSA 자동", "월 LTSA 제안 100↑",    NAVY),
        ("Phase 3", "확장",    "2027~2028",     "전 무상 선주 / 하이브리드 정식",       "LTSA 비중 20~30%",      ORANGE),
        ("Phase 4", "통합",    "2028~2029",     "HI-CMS 전 선박 / 예측정비 LTSA",      "LTSA 비중 35~40%",      ORANGE),
        ("Phase 5", "안정",    "2030",         "플릿 단위 LTSA 패키지",                "LTSA 비중 50%↑",        ORANGE),
    ]
    pw = Inches(2.45)
    gap = Inches(0.06)
    y = Inches(2.0)
    for i, (ph, name, period, task, kpi, color) in enumerate(phases):
        x = Inches(0.5) + i * (pw + gap)
        add_rect(s, x, y, pw, Inches(3.6), LIGHT_BLUE)
        add_rect(s, x, y, pw, Inches(0.55), color)
        add_text(s, x, y + Inches(0.07), pw, Inches(0.4),
                 ph, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), y + Inches(0.65),
                 pw - Inches(0.2), Inches(0.5), name,
                 size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), y + Inches(1.2),
                 pw - Inches(0.2), Inches(0.3), period,
                 size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        # divider
        add_rect(s, x + Inches(0.2), y + Inches(1.6), pw - Inches(0.4), Emu(1), LIGHT_GRAY)
        add_text(s, x + Inches(0.15), y + Inches(1.7),
                 pw - Inches(0.3), Inches(0.3), "핵심 과업",
                 size=9, color=GRAY_SUB, bold=True)
        add_text(s, x + Inches(0.15), y + Inches(1.95),
                 pw - Inches(0.3), Inches(1.1), task,
                 size=10, color=GRAY_TEXT)
        add_rect(s, x + Inches(0.2), y + Inches(3.0), pw - Inches(0.4), Emu(1), LIGHT_GRAY)
        add_text(s, x + Inches(0.15), y + Inches(3.05),
                 pw - Inches(0.3), Inches(0.3), "KPI",
                 size=9, color=GRAY_SUB, bold=True)
        add_text(s, x + Inches(0.15), y + Inches(3.27),
                 pw - Inches(0.3), Inches(0.3), kpi,
                 size=10, bold=True, color=NAVY)

    # gate timeline at bottom
    add_text(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.4),
             "단계별 의사결정 게이트 (Go / No-Go)",
             size=12, bold=True, color=NAVY)
    gates = [
        "G1 · AI 신뢰도 90%↑ + PM 만족도 4.0↑",
        "G2 · Closed Beta 갱신율 80%↑",
        "G3 · LTSA 비중 20%↑ + 매출 예측 ±15%↓",
        "G4 · LTSA 비중 35%↑",
    ]
    gw = Inches(3.05)
    gy = Inches(6.3)
    for i, g in enumerate(gates):
        x = Inches(0.5) + i * (gw + Inches(0.05))
        add_rect(s, x, gy, gw, Inches(0.6), WHITE, line=ORANGE)
        add_rect(s, x, gy, Inches(0.08), Inches(0.6), ORANGE)
        add_text(s, x + Inches(0.18), gy, gw - Inches(0.25), Inches(0.6),
                 g, size=10, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)


def slide_revenue_sim(prs):
    s = add_blank(prs)
    add_header(s, 13, "본론 5 · Q5 — 매출 시뮬레이션",
               "2026 → 2030: LTSA 비중 5% → 50%↑")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "보수 가정 — 적용 가능 모집단 5,900척, 단가 OEM 3사 LCA 70%, 선주 옵트인 단계 확대",
             size=12, color=GRAY_SUB)

    years = [
        ("2026", "50~100척",        0.05,  "5%",     "±25%"),
        ("2027", "500~800척",       0.13,  "10~15%", "±20%"),
        ("2028", "1,500~2,000척",   0.25,  "20~30%", "±15%"),
        ("2029", "2,500~3,500척",   0.38,  "35~40%", "±10%"),
        ("2030", "4,000~5,000척",   0.52,  "50%↑",   "±5%"),
    ]
    chart_left = Inches(2.5)
    chart_w = Inches(8.0)
    bar_h = Inches(0.5)
    gap = Inches(0.18)
    y0 = Inches(2.0)
    for i, (yr, ships, ratio, share, acc) in enumerate(years):
        y = y0 + i * (bar_h + gap)
        add_text(s, Inches(0.5), y + Inches(0.1), Inches(2.0), Inches(0.4),
                 yr, size=15, bold=True, color=NAVY)

        add_rect(s, chart_left, y + Inches(0.12), chart_w, Inches(0.28), LIGHT_BLUE)
        bar_w = Emu(int(chart_w * ratio))
        color = ORANGE if i == len(years) - 1 else NAVY
        add_rect(s, chart_left, y + Inches(0.12), bar_w, Inches(0.28), color)
        add_text(s, chart_left + bar_w + Inches(0.1), y + Inches(0.05),
                 Inches(2.0), Inches(0.4), share,
                 size=14, bold=True, color=color)

        add_text(s, Inches(10.6), y + Inches(0.0), Inches(2.4), Inches(0.3),
                 "가입 척수 (누적)", size=8, color=GRAY_SUB)
        add_text(s, Inches(10.6), y + Inches(0.22), Inches(2.4), Inches(0.3),
                 ships, size=10, bold=True, color=NAVY)

    add_text(s, Inches(2.5), Inches(5.55), Inches(8), Inches(0.3),
             "← 2026 5%                                                        2030  50%↑  →",
             size=9, color=GRAY_SUB, align=PP_ALIGN.CENTER)

    # KPI strip
    kpis = [
        ("9,890척",  "전체 모집단",         "HD현대 건조선 + AS 기존 선주"),
        ("5,900척",  "LTSA 적용 가능",     "전체의 약 60%"),
        ("70%",      "단가 가정",            "OEM 3사 LCA 평균 대비"),
        ("±5%",      "2030 매출 예측",     "OEM 3사 수준 안정성"),
    ]
    kw = Inches(3.05)
    ky = Inches(5.9)
    for i, (val, label, sub) in enumerate(kpis):
        x = Inches(0.5) + i * (kw + Inches(0.05))
        add_rect(s, x, ky, kw, Inches(0.78), LIGHT_BLUE)
        add_rect(s, x, ky, Inches(0.08), Inches(0.78), ORANGE)
        add_text(s, x + Inches(0.2), ky + Inches(0.04),
                 Inches(1.4), Inches(0.4), val,
                 size=18, bold=True, color=NAVY)
        add_text(s, x + Inches(1.45), ky + Inches(0.05),
                 kw - Inches(1.55), Inches(0.3), label,
                 size=10, bold=True, color=NAVY)
        add_text(s, x + Inches(1.45), ky + Inches(0.34),
                 kw - Inches(1.55), Inches(0.4), sub,
                 size=8, color=GRAY_SUB)

    add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.2),
             "※ Phase 1 PoC 종료 시점에 실측 데이터로 재보정 예정.",
             size=9, italic=True, color=GRAY_SUB, align=PP_ALIGN.CENTER)


def slide_before_after(prs):
    s = add_blank(prs)
    add_header(s, 14, "본론 6 — 업무 프로세스 개선 효과",
               "Before / After 종합 비교")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "1차 목표는 매출 증대(LTSA 50%↑) — 그 과정에서 R&R·KPI·데이터 활용 모두 자동 개선",
             size=12, color=GRAY_SUB)

    metrics = [
        ("PM 1척당 케이스 재구성",    "3~8일 (메일 수기)",       0.95,  "15~30분 (AI 자동)",      0.04),
        ("LTSA 제안 발송 지연",       "평균 2주~1개월",          1.00,  "실시간 (케이스 종료 직후)", 0.02),
        ("HI-CMS 활용률 (선주 노출)", "0% (사내 한정)",          0.0,   "100% (이상 감지 즉시)",  1.00),
        ("LTSA 매출 비중",            "5~10%",                  0.10,  "50%↑ (2030)",             1.00),
        ("매출 예측 정확도",          "±25% (단발성 의존)",      1.00,  "±5% (OEM 3사 수준)",      0.10),
        ("선주 이탈 리스크 측정",      "측정 불가",               0.0,   "갱신율·이탈율 실측 가능",   1.00),
    ]
    chart_top = Inches(1.95)
    row_h = Inches(0.55)
    gap_y = Inches(0.05)
    label_w = Inches(3.4)
    bar_zone_x = Inches(4.1)
    bar_zone_w = Inches(6.1)
    text_zone_x = bar_zone_x + bar_zone_w + Inches(0.1)
    text_zone_w = Inches(2.9)

    lg_y = Inches(1.45)
    add_rect(s, Inches(8.7), lg_y, Inches(0.3), Inches(0.18), GRAY_SUB)
    add_text(s, Inches(9.05), lg_y - Inches(0.02), Inches(1.0), Inches(0.25),
             "Before", size=10, color=GRAY_TEXT)
    add_rect(s, Inches(11.0), lg_y, Inches(0.3), Inches(0.18), ORANGE)
    add_text(s, Inches(11.35), lg_y - Inches(0.02), Inches(1.5), Inches(0.25),
             "After", size=10, color=GRAY_TEXT)

    for i, (label, b_text, b_pct, a_text, a_pct) in enumerate(metrics):
        y = chart_top + i * (row_h + gap_y)
        if i % 2 == 0:
            add_rect(s, Inches(0.5), y - Inches(0.02),
                     Inches(12.3), row_h + Inches(0.04), LIGHT_BLUE)
        add_text(s, Inches(0.6), y + Inches(0.13), label_w, Inches(0.3),
                 label, size=11, bold=True, color=NAVY)

        before_w = Emu(int(bar_zone_w * 0.45 * b_pct))
        add_rect(s, bar_zone_x, y + Inches(0.04),
                 max(before_w, Emu(1)), Inches(0.18), GRAY_SUB)
        after_w = Emu(int(bar_zone_w * a_pct))
        add_rect(s, bar_zone_x, y + Inches(0.27),
                 max(after_w, Emu(1)), Inches(0.20), ORANGE)
        add_text(s, text_zone_x, y + Inches(0.0), text_zone_w, Inches(0.25),
                 b_text, size=10, color=GRAY_SUB)
        add_text(s, text_zone_x, y + Inches(0.25), text_zone_w, Inches(0.25),
                 a_text, size=10, bold=True, color=ORANGE)

    add_rect(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.7), NAVY)
    add_rect(s, Inches(0.5), Inches(6.15), Inches(0.12), Inches(0.7), ORANGE)
    add_text(s, Inches(0.85), Inches(6.20), Inches(12), Inches(0.3),
             "한 줄 요약",
             size=10, bold=True, color=ORANGE)
    add_text(s, Inches(0.85), Inches(6.45), Inches(12), Inches(0.4),
             "PM 처리 능력 ×5~7   ·   LTSA 매출 비중 ×10   ·   매출 예측 정확도 5배 향상   ·   HI-CMS 활용률 0 → 100%",
             size=14, bold=True, color=WHITE)


def slide_conclusion(prs):
    s = add_blank(prs)
    add_header(s, 15, "결론", "4가지 메커니즘, 하나의 흐름")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "별도 영업 캠페인·신규 시장 진출·신규 인력 일괄 투입 없이, 이미 보유한 자산을 \"연결\"만 한다",
             size=12, color=GRAY_SUB)

    axes = [
        ("①", "무상 기술지원 포털", "입구",
         "선주가 자연스럽게 들어오는\n단일 창구",                        ORANGE),
        ("②", "AI 케이스 인계",     "통로 1",
         "PM이 메일을 뒤지지 않고\n즉시 LTSA 제안 가능",                  TEAL),
        ("③", "HI-CMS 통합 알림",   "통로 2",
         "선주가 안 와도 됐던 영역을\n\"와야 하는 영역\"으로",              NAVY),
        ("④", "LTSA 자연 전환",     "출구",
         "단발 매출 →\n안정적 구독+사건 하이브리드",                       ORANGE),
    ]
    cw = Inches(3.0)
    gap = Inches(0.1)
    y = Inches(2.0)
    for i, (idx, title, role, desc, color) in enumerate(axes):
        x = Inches(0.5) + i * (cw + gap)
        add_rect(s, x, y, cw, Inches(2.8), LIGHT_BLUE)
        add_rect(s, x, y, cw, Inches(0.5), color)
        add_text(s, x, y + Inches(0.05), cw, Inches(0.4),
                 role, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), y + Inches(0.7),
                 Inches(0.6), Inches(0.6), idx,
                 size=28, bold=True, color=color)
        add_text(s, x + Inches(0.85), y + Inches(0.75),
                 cw - Inches(1.0), Inches(0.7), title,
                 size=14, bold=True, color=NAVY)
        add_rect(s, x + Inches(0.3), y + Inches(1.55),
                 cw - Inches(0.6), Emu(1), LIGHT_GRAY)
        add_text(s, x + Inches(0.2), y + Inches(1.7),
                 cw - Inches(0.4), Inches(1.0), desc,
                 size=11, color=GRAY_TEXT)

        if i < len(axes) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + cw, y + Inches(1.3), gap, Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()

    # promise box
    add_rect(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.85), NAVY)
    add_rect(s, Inches(0.5), Inches(5.1), Inches(0.15), Inches(1.85), ORANGE)
    add_text(s, Inches(0.85), Inches(5.2), Inches(12), Inches(0.3),
             "단 한 줄의 약속", size=11, bold=True, color=ORANGE)
    add_text(s, Inches(0.85), Inches(5.5), Inches(12), Inches(0.7),
             '"메일을 열던 그 손이, 자연스럽게 LTSA 제안을 클릭하는 순간을 만든다."',
             size=22, bold=True, italic=True, color=WHITE)
    add_text(s, Inches(0.85), Inches(6.25), Inches(12), Inches(0.3),
             "회사 주제 두 가지 — 중장기 매출 증대 + 업무 프로세스 개선 — 동시 달성.",
             size=12, color=LIGHT_GRAY)
    add_text(s, Inches(0.85), Inches(6.55), Inches(12), Inches(0.3),
             "→ 2030년: LTSA 비중 50%↑  ·  매출 예측 정확도 ±5%  ·  PM 처리 능력 ×5~7",
             size=12, bold=True, color=ORANGE)

    # KPI line moved into the navy callout (slide 15) — drop the separate bottom line
    pass


def slide_appendix_qa(prs):
    s = add_blank(prs)
    add_header(s, 16, "부록 — 발표 시간 + 핵심 질문 답변",
               "10분 발표 구간별 메시지 + Q1~Q5 1줄 답변")

    # left: time allocation
    add_rect(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.5), NAVY)
    add_text(s, Inches(0.7), Inches(1.55), Inches(5.8), Inches(0.4),
             "발표 시간 배분 (10분)",
             size=12, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    times = [
        ("0:00~1:00",  "서론 — 두 주제 통합"),
        ("1:00~3:00",  "본론 1 — 단절·PM 부담"),
        ("3:00~4:30",  "본론 2 — LTSA 모델"),
        ("4:30~6:30",  "본론 3 — AI 인계 효과"),
        ("6:30~7:30",  "본론 4 — HI-CMS 알림"),
        ("7:30~9:00",  "본론 5 — 2030 로드맵"),
        ("9:00~10:00", "결론 + Q&A"),
    ]
    for i, (t, msg) in enumerate(times):
        y = Inches(2.05) + i * Inches(0.55)
        if i % 2 == 0:
            add_rect(s, Inches(0.5), y, Inches(6.0), Inches(0.5), LIGHT_BLUE)
        add_text(s, Inches(0.7), y, Inches(1.5), Inches(0.5),
                 t, size=10, bold=True, color=ORANGE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.3), y, Inches(4.0), Inches(0.5),
                 msg, size=11, color=GRAY_TEXT,
                 anchor=MSO_ANCHOR.MIDDLE)

    # right: Q1-Q5 answers
    add_rect(s, Inches(6.7), Inches(1.5), Inches(6.1), Inches(0.5), NAVY)
    add_text(s, Inches(6.9), Inches(1.55), Inches(5.9), Inches(0.4),
             "5개 질문 1줄 답변",
             size=12, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    answers = [
        ("Q1", "왜 LTSA?",          "OEM 3사가 매출 안정성 확보한 유일 구조"),
        ("Q2", "단절 해소?",        "무상 포털 자체가 LTSA 영업 도구"),
        ("Q3", "AI 효과?",          "3~8일 → 15~30분 (99% 시간 단축)"),
        ("Q4", "HI-CMS 역할?",     "이상→웹 알림→케이스→LTSA (클릭 1회)"),
        ("Q5", "2030 도달?",        "5-Phase / LTSA 비중 5% → 50%↑"),
    ]
    for i, (qk, qt, ans) in enumerate(answers):
        y = Inches(2.05) + i * Inches(0.78)
        add_rect(s, Inches(6.7), y, Inches(6.1), Inches(0.7), LIGHT_BLUE)
        add_rect(s, Inches(6.7), y, Inches(0.7), Inches(0.7), NAVY)
        add_text(s, Inches(6.7), y, Inches(0.7), Inches(0.7),
                 qk, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.5), y + Inches(0.05),
                 Inches(5.2), Inches(0.3), qt,
                 size=11, bold=True, color=NAVY)
        add_text(s, Inches(7.5), y + Inches(0.32),
                 Inches(5.2), Inches(0.4), ans,
                 size=10, color=GRAY_TEXT)


def slide_qna(prs):
    s = add_blank(prs)
    add_header(s, 17, "부록 — 예상 Q&A",
               "임원 질의 대응 6선")

    qas = [
        ("AI 신뢰도 부족 시 대응?",
         "PM 검토 5~15분 유지 / 신뢰도 90%↑ 도달까지 PoC 검증."),
        ("HI-CMS 데이터 보안·동의?",
         "선주별 옵트인 / OEM 3사 LCA 표준 약관 동등 수준 적용."),
        ("구독형이 선주 부담 키울 우려?",
         "하이브리드(기본료+사용량) 권고 — 저사용 선주 손해 없음."),
        ("무상→유상 전환 시 선주 저항?",
         "무상 그대로 유지 / 유상은 추가 가치(예측·우선)에만 한정."),
        ("필요 인력·예산 규모?",
         "Phase 1 기준 AI/플랫폼 PoC 5~7명 / 신규 영업 인력 불필요."),
        ("OEM 3사 견제 가능성?",
         "OEM은 자기 장비만 / HMS는 HD현대 건조선 전체 통합 — 차별점."),
    ]
    cw = Inches(6.05)
    rh = Inches(1.5)
    gap = Inches(0.15)
    y0 = Inches(1.7)
    for i, (q, a) in enumerate(qas):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * (cw + Inches(0.2))
        y = y0 + row * (rh + gap)
        add_rect(s, x, y, cw, rh, LIGHT_BLUE)
        add_rect(s, x, y, Inches(0.12), rh, ORANGE)
        add_text(s, x + Inches(0.3), y + Inches(0.1),
                 Inches(0.6), Inches(0.4), f"Q{i+1}",
                 size=14, bold=True, color=ORANGE)
        add_text(s, x + Inches(1.0), y + Inches(0.15),
                 cw - Inches(1.2), Inches(0.4), q,
                 size=12, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(0.7),
                 cw - Inches(0.5), Inches(0.7), "→ " + a,
                 size=11, color=GRAY_TEXT)


def slide_final(prs):
    s = add_blank(prs)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), NAVY)
    add_rect(s, Inches(0), Inches(7.4), Inches(13.33), Inches(0.1), ORANGE)

    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             "FINAL MESSAGE", size=13, bold=True, color=ORANGE)
    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.0),
             "메일을 열던 그 손이,",
             size=42, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(1.0),
             "자연스럽게 LTSA 제안을 클릭하는 순간.",
             size=42, bold=True, color=ORANGE)

    add_rect(s, Inches(0.7), Inches(4.5), Inches(0.5), Inches(0.06), ORANGE)
    add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.5),
             "무상 기술지원 → AI 케이스 인계 → HI-CMS 알림 → LTSA 자연 전환",
             size=18, color=WHITE)
    add_text(s, Inches(0.7), Inches(5.3), Inches(12), Inches(0.5),
             "2030년: LTSA 비중 50%↑  ·  매출 예측 정확도 ±5%  ·  OEM 3사와 동등한 안정 매출 구조",
             size=14, color=LIGHT_GRAY)

    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
             "감사합니다  ·  Thank You",
             size=20, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3),
             "HD현대마린솔루션테크",
             size=12, color=ORANGE)


# ── main ----------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_title,           # 1
        slide_executive,       # 2
        slide_toc,             # 3
        slide_intro,           # 4
        slide_status,          # 5
        slide_revenue_gap,     # 6
        slide_ltsa_flow,       # 7
        slide_ltsa_models,     # 8
        slide_ai_flow,         # 9
        slide_ai_impact,       # 10
        slide_hicms,           # 11
        slide_roadmap,         # 12
        slide_revenue_sim,     # 13
        slide_before_after,    # 14
        slide_conclusion,      # 15
        slide_appendix_qa,     # 16
        slide_qna,             # 17
        slide_final,           # 18
    ]

    global TOTAL_PAGES
    TOTAL_PAGES = len(builders)

    for build in builders:
        build(prs)

    prs.save(str(OUTPUT))
    print(f"✓ saved {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
